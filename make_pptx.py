"""Generate ISLA Bot presentation as a .pptx file."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── colour palette ──────────────────────────────────────────────────────────
BG      = RGBColor(0x0d, 0x11, 0x17)
SURFACE = RGBColor(0x16, 0x1b, 0x22)
BORDER  = RGBColor(0x30, 0x36, 0x3d)
TEXT    = RGBColor(0xe6, 0xed, 0xf3)
MUTED   = RGBColor(0x8b, 0x94, 0x9e)
GREEN   = RGBColor(0x3f, 0xb9, 0x50)
RED     = RGBColor(0xf8, 0x51, 0x49)
BLUE    = RGBColor(0x58, 0xa6, 0xff)
YELLOW  = RGBColor(0xd2, 0xa5, 0x20)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

# ── slide size 16:9 ─────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── helpers ─────────────────────────────────────────────────────────────────

def add_slide():
    sl = prs.slides.add_slide(blank_layout)
    # dark background
    bg = sl.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return sl


def box(sl, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0.75), radius=None):
    """Add a filled/bordered rectangle shape."""
    shape = sl.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.width = border_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
    else:
        shape.line.fill.background()
    return shape


def txt(sl, text, left, top, width, height,
        size=Pt(12), bold=False, italic=False, color=TEXT,
        align=PP_ALIGN.LEFT, wrap=True, valign=None):
    """Add a textbox."""
    txBox = sl.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if valign:
        tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return txBox


def txt_lines(sl, lines, left, top, width, height,
              size=Pt(11), color=TEXT, bold=False, line_spacing=Pt(16)):
    """Add a textbox with multiple paragraph lines (list of (text, color, bold, size))."""
    txBox = sl.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            item = (item, color, bold, size)
        t, c, b, s = item
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = t
        run.font.size = s
        run.font.bold = b
        run.font.color.rgb = c
        run.font.name = "Segoe UI"
    return txBox


def heading(sl, title, accent=GREEN):
    """Draw slide heading with green underline bar."""
    # title text
    txt(sl, title,
        Inches(0.35), Inches(0.18), Inches(12.6), Inches(0.55),
        size=Pt(22), bold=True, color=TEXT)
    # underline bar
    box(sl, Inches(0.35), Inches(0.72), Inches(12.6), Pt(3),
        fill_color=accent, border_color=None)


def bullet_lines(sl, items, left, top, width, height, size=Pt(11)):
    """Render bulleted list items."""
    txBox = sl.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = "▸  " + item
        run.font.size = size
        run.font.color.rgb = TEXT
        run.font.name = "Segoe UI"
    return txBox


def card(sl, left, top, width, height, label, value, desc="", val_color=GREEN):
    box(sl, left, top, width, height, fill_color=SURFACE, border_color=BORDER)
    txt(sl, label, left + Inches(0.1), top + Inches(0.07), width - Inches(0.2), Inches(0.25),
        size=Pt(7.5), bold=True, color=MUTED)
    txt(sl, value, left + Inches(0.1), top + Inches(0.28), width - Inches(0.2), Inches(0.55),
        size=Pt(22), bold=True, color=val_color)
    if desc:
        txt(sl, desc, left + Inches(0.1), top + Inches(0.82), width - Inches(0.2), Inches(0.3),
            size=Pt(8), color=MUTED)


def key_msg(sl, text, color=GREEN):
    """Green bar at bottom with key message."""
    bar_top = H - Inches(0.42)
    box(sl, 0, bar_top, W, Inches(0.42),
        fill_color=RGBColor(0x0a, 0x25, 0x10) if color == GREEN else RGBColor(0x08, 0x18, 0x2e),
        border_color=None)
    # top border line
    box(sl, 0, bar_top, W, Pt(1.5), fill_color=color, border_color=None)
    txt(sl, "💡  " + text,
        Inches(0.35), bar_top + Inches(0.07), Inches(12.6), Inches(0.3),
        size=Pt(10), color=color)


def callout(sl, text, left, top, width, height, color=GREEN):
    accent_colors = {GREEN: RGBColor(0x0a, 0x25, 0x10),
                     BLUE:  RGBColor(0x08, 0x18, 0x2e),
                     YELLOW:RGBColor(0x28, 0x20, 0x05),
                     RED:   RGBColor(0x28, 0x06, 0x05)}
    bg_c = accent_colors.get(color, RGBColor(0x0a, 0x25, 0x10))
    box(sl, left, top, width, height, fill_color=bg_c, border_color=None)
    box(sl, left, top, Pt(4), height, fill_color=color, border_color=None)
    txt(sl, text, left + Inches(0.12), top + Inches(0.08),
        width - Inches(0.18), height - Inches(0.12),
        size=Pt(10), color=TEXT, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()

# gradient-ish logo box
logo_w = Inches(5)
logo_h = Inches(1.1)
logo_l = (W - logo_w) / 2
logo_t = Inches(0.9)
lb = box(sl, logo_l, logo_t, logo_w, logo_h, fill_color=RGBColor(0x1a, 0x5c, 0x30), border_color=None)
txt(sl, "ISLA Bot", logo_l, logo_t, logo_w, logo_h,
    size=Pt(46), bold=True, color=RGBColor(0x0d, 0x11, 0x17), align=PP_ALIGN.CENTER)

txt(sl, "Trading Algorítmico com Machine Learning",
    Inches(1), Inches(2.15), Inches(11.33), Inches(0.45),
    size=Pt(17), bold=True, color=BLUE, align=PP_ALIGN.CENTER)

txt(sl, "BTC/USD · Paper Trading · Alpaca API · Random Forest · Sentimento LLM",
    Inches(1), Inches(2.62), Inches(11.33), Inches(0.35),
    size=Pt(11), color=MUTED, align=PP_ALIGN.CENTER)

# group box
gbox_w = Inches(6)
gbox_l = (W - gbox_w) / 2
box(sl, gbox_l, Inches(3.1), gbox_w, Inches(1.35), fill_color=SURFACE, border_color=BORDER)
txt(sl, "GRUPO", gbox_l, Inches(3.18), gbox_w, Inches(0.25),
    size=Pt(8), bold=True, color=MUTED, align=PP_ALIGN.CENTER)
txt(sl, "Raphael Malburg   ·   André Neves   ·   Vasco   ·   Beatriz Ferreira",
    gbox_l, Inches(3.45), gbox_w, Inches(0.35),
    size=Pt(12), bold=True, color=TEXT, align=PP_ALIGN.CENTER)

txt(sl, "Engenharia de Software e IA · Março 2026",
    Inches(1), Inches(4.6), Inches(11.33), Inches(0.3),
    size=Pt(10), color=MUTED, align=PP_ALIGN.CENTER)
txt(sl, "⚠  Projeto académico · Não é recomendação de investimento · Paper trading apenas",
    Inches(1), Inches(4.95), Inches(11.33), Inches(0.3),
    size=Pt(8.5), color=RGBColor(0x55, 0x5c, 0x64), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Agenda")

agenda = [
    ("1", "Motivação e Problema"),
    ("2", "Solução — O que é o ISLA Bot"),
    ("3", "Arquitetura do Sistema"),
    ("4", "Pipeline de Decisão (8 passos)"),
    ("5", "Modelo de Machine Learning"),
    ("6", "Métricas do Modelo"),
    ("7", "Gestão de Risco"),
    ("8", "Dashboard Web"),
    ("9", "Resultados — Backtest 6 Meses"),
    ("10", "Kelly vs Fixed 5%"),
    ("11", "Conclusões"),
    ("12", "Metodologia & Equipa"),
    ("13", "Próximos Passos"),
    ("14", "Referências"),
]

cols = 2
rows = (len(agenda) + 1) // 2
cell_w = Inches(6.1)
cell_h = Inches(0.44)
gap = Inches(0.1)
start_l = Inches(0.35)
start_t = Inches(0.92)

for i, (num, label) in enumerate(agenda):
    col = i % cols
    row = i // cols
    l = start_l + col * (cell_w + Inches(0.15))
    t = start_t + row * (cell_h + gap)
    box(sl, l, t, cell_w, cell_h, fill_color=SURFACE, border_color=BORDER)
    txt(sl, num, l + Inches(0.12), t + Inches(0.06), Inches(0.35), cell_h - Inches(0.1),
        size=Pt(16), bold=True, color=GREEN)
    txt(sl, label, l + Inches(0.5), t + Inches(0.09), cell_w - Inches(0.55), cell_h - Inches(0.1),
        size=Pt(11), color=TEXT)

key_msg(sl, "O bot resolve um problema real: mercados cripto operam 24/7, humanos não conseguem.", BLUE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — MOTIVAÇÃO
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Motivação e Problema")

stats = [
    ("24/7", RED, "Mercados cripto nunca fecham —\nmonitorização humana contínua é impossível"),
    ("≈70%", YELLOW, "dos traders individuais perdem dinheiro\npor decisões emocionais e inconsistentes"),
    ("ML", GREEN, "Dados históricos ricos permitem modelar\npadrões e aplicar regras sem emoção"),
]
stat_w = Inches(3.9)
stat_h = Inches(2.1)
stat_gap = Inches(0.2)
stat_t = Inches(0.95)
total_w = 3 * stat_w + 2 * stat_gap
stat_start = (W - total_w) / 2

for i, (num, col, desc) in enumerate(stats):
    l = stat_start + i * (stat_w + stat_gap)
    box(sl, l, stat_t, stat_w, stat_h, fill_color=SURFACE, border_color=BORDER)
    txt(sl, num, l, stat_t + Inches(0.2), stat_w, Inches(0.85),
        size=Pt(46), bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(sl, desc, l + Inches(0.15), stat_t + Inches(1.1), stat_w - Inches(0.3), Inches(0.9),
        size=Pt(10), color=MUTED, align=PP_ALIGN.CENTER, wrap=True)

callout(sl,
        "Pergunta de investigação: É possível construir um sistema algorítmico que identifique "
        "oportunidades em BTC/USD e supere Buy & Hold usando indicadores técnicos e ML, "
        "com gestão de risco automática?",
        Inches(0.35), Inches(3.22), Inches(12.6), Inches(0.85), BLUE)

key_msg(sl, "A vantagem do bot não é ser mais inteligente — é ser consistente e disciplinado.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SOLUÇÃO
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "A Solução — ISLA Bot")

bullet_lines(sl, [
    "Bot de trading algorítmico para BTC/USD em paper trading",
    "Analisa candles de 1 hora com 17 indicadores técnicos",
    "Random Forest prevê se o TP será atingido antes do SL",
    "Gate de sentimento (Gemini Flash) bloqueia entradas em mercado muito bearish",
    "Gestão de risco automática: SL/TP baseados em ATR, circuit breaker, exit watchdog a cada 10s",
    "Dashboard web em tempo real com gráficos interativos",
], Inches(0.35), Inches(0.92), Inches(6.6), Inches(3.0))

# right column cards
right_l = Inches(7.2)
card(sl, right_l, Inches(0.92), Inches(5.8), Inches(1.0),
     "Exchange", "Alpaca Markets", "Paper trading, sem risco real", TEXT)

box(sl, right_l, Inches(2.02), Inches(5.8), Inches(1.2), fill_color=SURFACE, border_color=BORDER)
txt(sl, "Stack", right_l + Inches(0.1), Inches(2.08), Inches(5.6), Inches(0.22),
    size=Pt(7.5), bold=True, color=MUTED)
tags = [("Python", BLUE), ("Flask", TEXT), ("Scikit-learn", BLUE),
        ("SQLite", TEXT), ("Plotly.js", TEXT), ("Gemini", YELLOW)]
tag_l = right_l + Inches(0.1)
tag_t = Inches(2.32)
for tag, tc in tags:
    tw = Inches(0.9)
    box(sl, tag_l, tag_t, tw, Inches(0.22), fill_color=RGBColor(0x20, 0x28, 0x34), border_color=None)
    txt(sl, tag, tag_l + Inches(0.05), tag_t, tw - Inches(0.05), Inches(0.22), size=Pt(8), color=tc)
    tag_l += tw + Inches(0.08)
    if tag_l > right_l + Inches(5.5):
        tag_l = right_l + Inches(0.1)
        tag_t += Inches(0.28)

card(sl, right_l, Inches(3.32), Inches(5.8), Inches(0.75),
     "Timeframe · Ativo", "1H candles · BTC/USD", "", TEXT)

key_msg(sl, "Pipeline end-to-end: dados de mercado → ML → sentimento → risco → ordem na exchange.", BLUE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ARQUITETURA
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Arquitetura do Sistema")

arch_items = [
    # (label, sub, color, row, col_span)
    [("Alpaca API", "Candles 1H · Ordens · Conta", BLUE),
     ("OpenRouter / Gemini", "Análise de sentimento LLM", YELLOW)],
    [("src/data.py", "Fetch + 17 indicadores", TEXT),
     ("src/model.py", "Random Forest · predict_proba", BLUE),
     ("src/strategy.py", "Sizing · SL/TP · Checks", TEXT),
     ("src/broker.py", "Alpaca · ordens · posições", GREEN)],
    [("src/app.py — Flask", "bot_loop (1H) · exit_watchdog (10s) · REST API · Threads", BLUE),
     ("SQLite", "trades · runs · equity", TEXT)],
    [("Dashboard Web", "Jinja2 · Plotly.js · Polling 15s", GREEN)],
]

row_tops  = [Inches(0.95), Inches(1.82), Inches(2.7), Inches(3.58)]
row_h     = Inches(0.65)
arrow_gap = Inches(0.2)
total_arch_w = Inches(12.6)
arch_l = Inches(0.35)

for ri, row_items in enumerate(arch_items):
    n = len(row_items)
    item_w = (total_arch_w - arrow_gap * (n - 1)) / n
    t = row_tops[ri]
    for ci, (label, sub, col) in enumerate(row_items):
        l = arch_l + ci * (item_w + arrow_gap)
        border = col if col != TEXT else BORDER
        box(sl, l, t, item_w, row_h, fill_color=SURFACE, border_color=border)
        txt(sl, label, l + Inches(0.1), t + Inches(0.07), item_w - Inches(0.15), Inches(0.28),
            size=Pt(10), bold=True, color=col)
        txt(sl, sub, l + Inches(0.1), t + Inches(0.33), item_w - Inches(0.15), Inches(0.28),
            size=Pt(8.5), color=MUTED)
    # arrow between rows
    if ri < len(arch_items) - 1:
        arr_t = row_tops[ri] + row_h + Inches(0.02)
        txt(sl, "↓", Inches(6.5), arr_t, Inches(0.4), Inches(0.18),
            size=Pt(12), color=BORDER, align=PP_ALIGN.CENTER)

key_msg(sl, "Dois threads em background: bot_loop (1H) e exit_watchdog (10s) — o servidor Flask nunca bloqueia.", BLUE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PIPELINE
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Pipeline de Decisão — A cada candle de 1H")

analysis = [
    ("1", "Fetch Market Data", "90 dias · 1H · ~2.160 barras BTC/USD via Alpaca", MUTED),
    ("2", "Feature Engineering", "RSI, MACD, EMA-20/50, BB, ATR, OBV + 10 mais (17 total)", MUTED),
    ("3", "ML Prediction — Random Forest", "P(TP antes SL em 12H) ≥ 55% → LONG · senão → FLAT", BLUE),
    ("4", "Drift Detection", "Confiança média últimas 30 runs < 50% → aviso", MUTED),
    ("5", "Sentiment Gate + Circuit Breaker", "Gemini ≤ −0.5 bloqueia · perda diária >10% para trading", MUTED),
]
execution = [
    ("6", "Position Sizing", "equity × 5% ÷ stop_pct · onde stop_pct = 1×ATR ÷ preço", MUTED),
    ("7", "Market BUY + Disaster Stop", "Ordem de mercado + stop-limit a 3×ATR (safety net de crash)", GREEN),
    ("8", "Exit Watchdog — a cada 10s", "SL atingido (1×ATR) · TP atingido (2×ATR) · max 12H · sinal flip", GREEN),
]

def pipe_section(sl, items, left, top, width, label, label_color):
    txt(sl, label, left + Inches(0.06), top, width, Inches(0.22),
        size=Pt(8.5), bold=True, color=label_color)
    box(sl, left, top, Pt(3), Inches(0.22), fill_color=label_color, border_color=None)
    t = top + Inches(0.26)
    for num, title, desc, nc in items:
        row_h = Inches(0.58)
        box(sl, left, t, width, row_h, fill_color=SURFACE,
            border_color=nc if nc != MUTED else BORDER)
        # number bubble
        box(sl, left + Inches(0.1), t + Inches(0.15), Inches(0.28), Inches(0.28),
            fill_color=RGBColor(0x20, 0x28, 0x34), border_color=None)
        txt(sl, num, left + Inches(0.1), t + Inches(0.14), Inches(0.28), Inches(0.3),
            size=Pt(8.5), bold=True, color=nc if nc != MUTED else MUTED, align=PP_ALIGN.CENTER)
        txt(sl, title, left + Inches(0.47), t + Inches(0.07), width - Inches(0.55), Inches(0.24),
            size=Pt(10), bold=True, color=TEXT)
        txt(sl, desc, left + Inches(0.47), t + Inches(0.3), width - Inches(0.55), Inches(0.22),
            size=Pt(8.5), color=MUTED)
        t += row_h + Inches(0.05)
    return t

col_w = Inches(6.1)
pipe_section(sl, analysis, Inches(0.35), Inches(0.92), col_w, "Análise (1–5)", BLUE)
end_t = pipe_section(sl, execution, Inches(6.58), Inches(0.92), col_w, "Execução (6–8)", GREEN)

callout(sl,
        "O watchdog verifica SL e TP a cada 10 segundos — não depende da exchange, "
        "o que elimina ordens órfãs e permite saídas baseadas em qualquer condição.",
        Inches(6.58), end_t + Inches(0.05), col_w, Inches(0.6), GREEN)

key_msg(sl, "O bot entra pouco e sai rápido — selectividade é a principal fonte de vantagem.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MODELO ML
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Modelo de Machine Learning")

col_w = Inches(6.1)

# left col
txt(sl, "Objetivo do modelo", Inches(0.35), Inches(0.97), col_w, Inches(0.28),
    size=Pt(12), bold=True, color=BLUE)
callout(sl,
        '"O Take-Profit (2×ATR) será atingido antes do Stop-Loss (1×ATR) nos próximos 12 candles?"',
        Inches(0.35), Inches(1.3), col_w, Inches(0.75), GREEN)
txt(sl, "Target directamente alinhado com a lógica real de saída.\nPrecisão = taxa de vitória nas posições abertas.",
    Inches(0.35), Inches(2.1), col_w, Inches(0.45), size=Pt(9.5), color=MUTED)

txt(sl, "Algoritmo & Validação", Inches(0.35), Inches(2.65), col_w, Inches(0.28),
    size=Pt(12), bold=True, color=BLUE)
bullet_lines(sl, [
    "Random Forest · 200 árvores · profundidade máx. 8",
    "StandardScaler em todos os features",
    "Walk-forward validation — TimeSeriesSplit 5 folds",
    "Sem data leakage temporal",
], Inches(0.35), Inches(2.98), col_w, Inches(1.2))

# right col — feature importance table
right_l = Inches(6.8)
txt(sl, "Top features (importância)", right_l, Inches(0.97), col_w, Inches(0.28),
    size=Pt(12), bold=True, color=BLUE)

headers = [("#", Inches(0.4)), ("Feature", Inches(1.8)), ("Peso", Inches(1.2))]
rows_data = [
    ("1", "OBV", "11.9%", GREEN),
    ("2", "ATR-14", "10.0%", GREEN),
    ("3", "EMA-20", "9.6%", GREEN),
    ("4", "EMA-50", "8.9%", TEXT),
    ("5", "BB High / Low", "~8.2%", TEXT),
    ("6", "MACD + Signal", "5.2%", TEXT),
    ("7", "RSI-14", "5.0%", MUTED),
]
th = Inches(0.28)
row_h = Inches(0.35)
table_t = Inches(1.3)
# header row
box(sl, right_l, table_t, col_w, th, fill_color=SURFACE, border_color=BORDER)
col_offsets = [Inches(0.05), Inches(0.55), Inches(4.2)]
col_labels  = ["#", "Feature", "Peso"]
for j, (lbl, co) in enumerate(zip(col_labels, col_offsets)):
    txt(sl, lbl, right_l + co, table_t + Inches(0.04), Inches(1.5), th,
        size=Pt(8), bold=True, color=MUTED)

for i, (n, feat, pct, pc) in enumerate(rows_data):
    rt = table_t + th + i * row_h
    bg = RGBColor(0x0d, 0x11, 0x17) if i % 2 == 0 else SURFACE
    box(sl, right_l, rt, col_w, row_h, fill_color=bg, border_color=BORDER)
    for j, (val, co) in enumerate(zip([n, feat, pct], col_offsets)):
        c = pc if j == 2 else TEXT
        txt(sl, val, right_l + co, rt + Inches(0.06), Inches(2.5), row_h,
            size=Pt(9.5), color=c)

callout(sl,
        "Volume (OBV) e volatilidade (ATR) superam indicadores de momentum — onde há volume há movimento.",
        right_l, table_t + th + len(rows_data) * row_h + Inches(0.08), col_w, Inches(0.62), BLUE)

key_msg(sl, "Walk-forward garante que cada previsão usa apenas dados do passado — como numa operação real.", BLUE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — MÉTRICAS ML
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Métricas do Modelo (holdout)")

metrics = [
    ("Accuracy",  "67.9%", GREEN,  "previsões corretas"),
    ("Precision", "23.6%", YELLOW, "taxa de vitória real"),
    ("Recall",    "10.9%", YELLOW, "oportunidades captadas"),
    ("F1 Score",  "0.149", MUTED,  "P × R balanceados"),
]
mw = Inches(3.0)
mg = Inches(0.2)
mt = Inches(0.95)
ml = (W - 4 * mw - 3 * mg) / 2
for i, (lbl, val, col, desc) in enumerate(metrics):
    card(sl, ml + i * (mw + mg), mt, mw, Inches(1.1), lbl, val, desc, col)

col_w = Inches(6.1)
txt(sl, "O que significam estes números", Inches(0.35), Inches(2.3), col_w, Inches(0.28),
    size=Pt(12), bold=True, color=BLUE)
bullet_lines(sl, [
    "Precision 23.6% — o bot ganha 1 em cada 4 trades que abre",
    "Break-even para RR 2:1 é 33.3% — estamos abaixo",
    "Alta accuracy (68%) porque o modelo acerta a maior parte dos FLATs",
    "Recall baixo = poucas entradas — o bot é muito selectivo",
], Inches(0.35), Inches(2.62), col_w, Inches(1.5))

right_l = Inches(6.8)
txt(sl, "Dados de treino", right_l, Inches(2.3), col_w, Inches(0.28),
    size=Pt(12), bold=True, color=BLUE)
train_rows = [
    ("Amostras treino", "6.967", GREEN),
    ("Amostras teste",  "1.742", TEXT),
    ("Features",        "17", TEXT),
    ("Target",          "TP antes SL em 12H", TEXT),
    ("Validação",       "Walk-forward 5 folds", TEXT),
]
rh = Inches(0.35)
for i, (k, v, vc) in enumerate(train_rows):
    rt = Inches(2.62) + i * rh
    bg = SURFACE if i % 2 == 0 else BG
    box(sl, right_l, rt, col_w, rh, fill_color=bg, border_color=BORDER)
    txt(sl, k, right_l + Inches(0.1), rt + Inches(0.07), Inches(2.8), rh, size=Pt(10), color=TEXT)
    txt(sl, v, right_l + Inches(3.0), rt + Inches(0.07), Inches(2.9), rh, size=Pt(10), color=vc, bold=(vc==GREEN))

key_msg(sl, "Precision abaixo de 33.3% significa edge negativo por trade — o ganho vem da selectividade, não da precisão.", YELLOW)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — GESTÃO DE RISCO
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Gestão de Risco")

risk_items = [
    ("🔴🟢", "Stop-Loss (1×ATR) & Take-Profit (2×ATR)",
     "Rácio RR de 2:1. Geridos pelo exit watchdog a cada 10 segundos — não por ordens na exchange. Uma posição vencedora compensa 2 perdedoras."),
    ("⚡", "Circuit Breaker + Max Hold (12H)",
     "Bloqueia novos trades se perda diária ultrapassar 10% do capital. Fecha posição automaticamente após 12 candles mesmo sem SL/TP."),
    ("⚠️", "Disaster Stop (3×ATR) na Alpaca",
     "Única ordem colocada na exchange — stop-limit a 3×ATR abaixo da entrada. Só dispara se o servidor cair. Protege contra perda catastrófica sem interferir com saídas normais."),
    ("📊", "Position Sizing (5% fixo)",
     "Risco fixo de 5% do capital por trade, ajustado pela volatilidade (ATR). Testámos Kelly Criterion — não melhora com edge negativo por trade."),
]

rw = Inches(6.1)
rh = Inches(2.2)
rg = Inches(0.15)
rt0 = Inches(0.92)
rl0 = Inches(0.35)

for i, (icon, title, desc) in enumerate(risk_items):
    col = i % 2
    row = i // 2
    l = rl0 + col * (rw + rg)
    t = rt0 + row * (rh + rg)
    box(sl, l, t, rw, rh, fill_color=SURFACE, border_color=BORDER)
    txt(sl, icon, l + Inches(0.12), t + Inches(0.1), Inches(0.5), Inches(0.45), size=Pt(18))
    txt(sl, title, l + Inches(0.12), t + Inches(0.52), rw - Inches(0.2), Inches(0.32),
        size=Pt(10.5), bold=True, color=TEXT)
    txt(sl, desc, l + Inches(0.12), t + Inches(0.85), rw - Inches(0.2), Inches(1.2),
        size=Pt(9), color=MUTED, wrap=True)

key_msg(sl, "Duas camadas de protecção: app-side (lógica real) + exchange-side (safety net) — resiliente a falhas.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — DASHBOARD
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Dashboard Web em Tempo Real")

col_w = Inches(6.1)
txt(sl, "Interface ao vivo", Inches(0.35), Inches(0.97), col_w, Inches(0.28),
    size=Pt(12), bold=True, color=BLUE)
bullet_lines(sl, [
    "Preço BTC e sinal do bot actualizados a cada 15 segundos sem page refresh",
    "Gráfico de velas interactivo com zoom, pan, scroll (estilo TradingView)",
    "Linhas de SL e TP sobrepostas no gráfico com etiqueta de preço ao vivo",
    "RSI-14 e MACD sincronizados com o gráfico principal",
    "Kill switch para fechar todas as posições imediatamente",
], Inches(0.35), Inches(1.3), col_w, Inches(1.6))

txt(sl, "Secções informativas", Inches(0.35), Inches(3.0), col_w, Inches(0.28),
    size=Pt(12), bold=True, color=BLUE)
bullet_lines(sl, [
    "📊  Showcase 6 meses — equity curve + log de 23 trades",
    "🤖  Métricas ML + gráfico de feature importance",
    "🔄  Pipeline explicado passo a passo com thresholds reais",
], Inches(0.35), Inches(3.32), col_w, Inches(0.9))

right_l = Inches(6.8)

# polling card
box(sl, right_l, Inches(0.97), col_w, Inches(1.05), fill_color=SURFACE, border_color=BORDER)
txt(sl, "Polling intervals", right_l + Inches(0.1), Inches(1.03), col_w - Inches(0.2), Inches(0.22),
    size=Pt(7.5), bold=True, color=MUTED)
txt(sl, "15s  →  preço + posição + SL/TP no gráfico", right_l + Inches(0.1), Inches(1.28),
    col_w - Inches(0.2), Inches(0.28), size=Pt(10), color=TEXT)
txt(sl, "30s  →  sinal ML + trades + métricas completas", right_l + Inches(0.1), Inches(1.6),
    col_w - Inches(0.2), Inches(0.28), size=Pt(10), color=TEXT)

# REST API card
box(sl, right_l, Inches(2.12), col_w, Inches(1.4), fill_color=SURFACE, border_color=BORDER)
txt(sl, "REST API", right_l + Inches(0.1), Inches(2.18), col_w - Inches(0.2), Inches(0.22),
    size=Pt(7.5), bold=True, color=MUTED)
api_lines = [
    ("GET   /api/live_stats", GREEN),
    ("GET   /api/dashboard_data", GREEN),
    ("GET   /api/showcase", GREEN),
    ("POST  /api/kill_switch", RED),
]
for i, (line, lc) in enumerate(api_lines):
    txt(sl, line, right_l + Inches(0.1), Inches(2.43) + i * Inches(0.28),
        col_w - Inches(0.2), Inches(0.28), size=Pt(9.5), color=lc)

callout(sl,
        "Sem WebSockets — polling assíncrono com fetch(). Plotly.js actualiza gráficos in-place preservando o zoom do utilizador.",
        right_l, Inches(3.62), col_w, Inches(0.6), BLUE)

key_msg(sl, "O dashboard é a 'cabine de pilotagem' — tudo visível numa única página sem necessidade de refresh.", BLUE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — RESULTADOS
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Resultados — Backtest Walk-Forward 6 Meses")

# hero metric
txt(sl, "Outperformance vs Buy & Hold",
    Inches(0.5), Inches(1.0), Inches(5.5), Inches(0.3),
    size=Pt(10), color=MUTED, align=PP_ALIGN.CENTER)
txt(sl, "+10.25%",
    Inches(0.5), Inches(1.3), Inches(5.5), Inches(1.4),
    size=Pt(64), bold=True, color=GREEN, align=PP_ALIGN.CENTER)
txt(sl, "estratégia −0.50%  ·  buy & hold −10.74%",
    Inches(0.5), Inches(2.75), Inches(5.5), Inches(0.3),
    size=Pt(10), color=MUTED, align=PP_ALIGN.CENTER)

# mini cards
mini_w = Inches(2.5)
mini_h = Inches(0.9)
mini_gap = Inches(0.3)
mini_l = Inches(0.5) + (Inches(5.5) - 2*mini_w - mini_gap) / 2
card(sl, mini_l, Inches(3.15), mini_w, mini_h, "Win Rate", "52.2%", "", GREEN)
card(sl, mini_l + mini_w + mini_gap, Inches(3.15), mini_w, mini_h, "Trades", "23", "", TEXT)

# right: table + callout
right_l = Inches(6.8)
col_w = Inches(6.1)

result_rows = [
    ("Max Drawdown",  "−12.34%", RED),
    ("Profit Factor", "1.05", TEXT),
    ("Avg Win",       "+$1.23", GREEN),
    ("Avg Loss",      "−$1.28", RED),
    ("Sharpe Ratio",  "−0.05", YELLOW),
]
# table header
rh = Inches(0.3)
box(sl, right_l, Inches(0.97), col_w, rh, fill_color=SURFACE, border_color=BORDER)
txt(sl, "Métrica", right_l + Inches(0.1), Inches(0.97), Inches(3.5), rh, size=Pt(8.5), bold=True, color=MUTED)
txt(sl, "Valor",   right_l + Inches(3.7), Inches(0.97), Inches(2.2), rh, size=Pt(8.5), bold=True, color=MUTED)
for i, (k, v, vc) in enumerate(result_rows):
    t = Inches(0.97) + rh + i * rh
    bg = SURFACE if i % 2 else BG
    box(sl, right_l, t, col_w, rh, fill_color=bg, border_color=BORDER)
    txt(sl, k, right_l + Inches(0.1), t + Inches(0.05), Inches(3.5), rh, size=Pt(10), color=TEXT)
    txt(sl, v, right_l + Inches(3.7), t + Inches(0.05), Inches(2.2), rh, size=Pt(10), color=vc, bold=True)

callout(sl,
        "Contexto: período fortemente bearish para BTC (−10.7%). O bot preservou capital ao ficar maioritariamente FLAT — a vantagem vem de não perder, não de ganhar muito.",
        right_l, Inches(2.62), col_w, Inches(0.85), YELLOW)

key_msg(sl, "Em mercado em baixa, não perder já é ganhar — o bot superou o benchmark em +10.25%.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — KELLY
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Análise — Kelly Criterion vs Fixed 5%")

col_w = Inches(6.1)
# table
kelly_rows = [
    ("Fixed 5% (atual)", "+4.72%", "0.17", "22.1%", True),
    ("Full Kelly",        "+3.11%", "0.19", "10.5%", False),
    ("Half Kelly",        "+1.52%", "0.16", "5.7%",  False),
    ("Quarter Kelly",     "+0.79%", "0.16", "2.9%",  False),
]
rh = Inches(0.38)
th = Inches(0.3)
tt = Inches(0.95)
box(sl, Inches(0.35), tt, col_w, th, fill_color=SURFACE, border_color=BORDER)
for j, hdr in enumerate(["Estratégia", "Retorno", "Sharpe", "Max DD"]):
    offsets = [Inches(0.1), Inches(2.8), Inches(4.2), Inches(5.2)]
    txt(sl, hdr, Inches(0.35) + offsets[j], tt + Inches(0.05), Inches(1.5), th,
        size=Pt(8.5), bold=True, color=MUTED)
for i, (strat, ret, sh, dd, winner) in enumerate(kelly_rows):
    t = tt + th + i * rh
    bg = RGBColor(0x0a, 0x25, 0x10) if winner else (SURFACE if i % 2 else BG)
    bc = GREEN if winner else BORDER
    box(sl, Inches(0.35), t, col_w, rh, fill_color=bg, border_color=bc)
    vals = [strat, ret, sh, dd]
    colors = [TEXT, GREEN if ret.startswith("+") else YELLOW, TEXT, RED]
    offsets = [Inches(0.1), Inches(2.8), Inches(4.2), Inches(5.2)]
    for j, (v, c, o) in enumerate(zip(vals, colors, offsets)):
        txt(sl, v, Inches(0.35) + o, t + Inches(0.07), Inches(1.5), rh,
            size=Pt(10), color=c, bold=(winner and j == 0))

right_l = Inches(6.8)
txt(sl, "Porquê Kelly não funciona aqui", right_l, Inches(0.97), col_w, Inches(0.28),
    size=Pt(12), bold=True, color=BLUE)

callout(sl,
        "f* = (p × b − q) / b\n"
        "Com p=0.236, b=2  →  f* = −14.6% (não apostar)\n\n"
        "O modelo prevê confiança de ~55–65% mas a precisão real é 23.6%.\n"
        "Kelly usa probabilidades optimistas → subdimensiona ganhos.",
        right_l, Inches(1.3), col_w, Inches(1.35), YELLOW)

callout(sl,
        "Quando Kelly funcionaria: precisão ≥ 33.3% (break-even para RR 2:1). Abaixo disso, fixed fraction ganha sempre.",
        right_l, Inches(2.75), col_w, Inches(0.75), GREEN)

key_msg(sl, "Kelly só é válido com edge positivo por trade — primeiro melhorar o modelo, depois aplicar Kelly.", YELLOW)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — CONCLUSÕES
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Conclusões")

col_w = Inches(6.1)
txt(sl, "O que conseguimos", Inches(0.35), Inches(0.97), col_w, Inches(0.28),
    size=Pt(12), bold=True, color=BLUE)
bullet_lines(sl, [
    "+10.25% de outperformance vs Buy & Hold num período bearish",
    "Pipeline completo end-to-end com 8 camadas de decisão",
    "Sistema de exit management robusto — watchdog 10s + disaster stop",
    "Dashboard profissional com atualização em tempo real",
], Inches(0.35), Inches(1.3), col_w, Inches(1.4))

right_l = Inches(6.8)
txt(sl, "Limitações identificadas", right_l, Inches(0.97), col_w, Inches(0.28),
    size=Pt(12), bold=True, color=RED)
bullet_lines(sl, [
    "Precision 23.6% — abaixo do break-even (33.3%)",
    "Apenas 23 trades em 6 meses — baixa frequência de sinal",
    "Kelly Criterion não aplicável sem edge positivo por trade",
    "Paper trading — slippage real pode ser maior",
], right_l, Inches(1.3), col_w, Inches(1.4))

callout(sl,
        "O bot demonstra que selectividade inteligente — ficar FLAT na maioria das barras — pode superar Buy & Hold em mercados em baixa. "
        "O valor está na ausência de maus trades, não na qualidade de cada trade individual.",
        Inches(0.35), Inches(2.85), Inches(12.6), Inches(0.85), BLUE)

key_msg(sl, "Evitar perdas é tão valioso quanto ganhar — o bot prova que disciplina algorítmica supera intuição humana.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — METODOLOGIA & EQUIPA
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Metodologia & Contribuições da Equipa")

col_w = Inches(6.1)
txt(sl, "Processo de Desenvolvimento", Inches(0.35), Inches(0.97), col_w, Inches(0.28),
    size=Pt(12), bold=True, color=BLUE)
bullet_lines(sl, [
    "Abordagem iterativa — ciclos semanais de build → test → backtest → refine",
    "Controlo de versões via Git — commits atómicos por feature",
    "Testes de integração com dados reais da Alpaca (paper trading)",
    "Walk-forward backtest como critério de aceitação de cada modelo",
    "Deploy contínuo em Railway — dashboard sempre acessível ao grupo",
], Inches(0.35), Inches(1.3), col_w, Inches(1.6))

callout(sl,
        "Nenhuma funcionalidade considerada 'pronta' sem backtest validar métricas iguais ou superiores à iteração anterior.",
        Inches(0.35), Inches(3.0), col_w, Inches(0.6), GREEN)

right_l = Inches(6.8)
txt(sl, "Contribuições por Membro", right_l, Inches(0.97), col_w, Inches(0.28),
    size=Pt(12), bold=True, color=BLUE)

members = [
    ("Raphael Malburg", "Arquitectura do sistema · Flask API · Dashboard web · Deploy Railway · Integração Alpaca"),
    ("André Neves",     "Pipeline ML · Feature engineering · Backtesting walk-forward · Análise de métricas"),
    ("Vasco",           "Gestão de risco · Exit watchdog · Análise Kelly Criterion · Position sizing"),
    ("Beatriz Ferreira","Integração sentimento LLM · Testes de integração · Documentação · Análise de resultados"),
]
mh = Inches(0.72)
mg = Inches(0.08)
for i, (name, role) in enumerate(members):
    t = Inches(1.3) + i * (mh + mg)
    box(sl, right_l, t, col_w, mh, fill_color=SURFACE, border_color=BORDER)
    txt(sl, name, right_l + Inches(0.12), t + Inches(0.08), col_w - Inches(0.2), Inches(0.28),
        size=Pt(10.5), bold=True, color=TEXT)
    txt(sl, role, right_l + Inches(0.12), t + Inches(0.36), col_w - Inches(0.2), Inches(0.3),
        size=Pt(8.5), color=MUTED)

key_msg(sl, "Cada módulo foi desenvolvido e testado independentemente antes da integração — falhas identificadas cedo.", BLUE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — PRÓXIMOS PASSOS
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Próximos Passos")

col_w = Inches(6.1)
txt(sl, "Prioridade — melhorar o modelo", Inches(0.35), Inches(0.97), col_w, Inches(0.28),
    size=Pt(12), bold=True, color=BLUE)
bullet_lines(sl, [
    "Substituir Random Forest por XGBoost / LightGBM — melhor para séries temporais financeiras",
    "Aumentar threshold para ≥ 65% de confiança — menos trades, precision maior",
    "Adicionar funding rates e open interest como features",
], Inches(0.35), Inches(1.3), col_w, Inches(1.1))
callout(sl,
        "Meta: precision ≥ 33.3% → Kelly Criterion passa a funcionar → sistema tem edge estatístico positivo por trade.",
        Inches(0.35), Inches(2.5), col_w, Inches(0.72), GREEN)

right_l = Inches(6.8)
txt(sl, "Sistema & Infraestrutura", right_l, Inches(0.97), col_w, Inches(0.28),
    size=Pt(12), bold=True, color=BLUE)
bullet_lines(sl, [
    "Trailing stop — protege lucros em tendências prolongadas",
    "✓ Deploy em cloud (Railway) — operacional 24/7",
    "Alertas por Telegram/Discord quando trade é executado",
], right_l, Inches(1.3), col_w, Inches(1.1))
callout(sl,
        "Sistema já operacional em Railway com dashboard público. Próxima barreira é elevar precision acima do break-even para live trading.",
        right_l, Inches(2.5), col_w, Inches(0.72), BLUE)

key_msg(sl, "Um passo de cada vez: primeiro precision > 33%, depois Kelly, depois capital real.", BLUE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — REFERÊNCIAS
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Referências")

col_w = Inches(6.1)
txt(sl, "Algoritmos & ML", Inches(0.35), Inches(0.97), col_w, Inches(0.28),
    size=Pt(12), bold=True, color=BLUE)
bullet_lines(sl, [
    "Breiman, L. (2001). Random Forests. Machine Learning, 45(1), 5–32.",
    "Prado, M. L. (2018). Advances in Financial Machine Learning. Wiley.",
    "Chan, E. P. (2013). Algorithmic Trading: Winning Strategies. Wiley.",
    "Kelly, J. L. (1956). A New Interpretation of Information Rate. Bell System Technical Journal.",
    "Wilder, J. W. (1978). New Concepts in Technical Trading Systems. — base para ATR e RSI.",
], Inches(0.35), Inches(1.3), col_w, Inches(2.2))

right_l = Inches(6.8)
txt(sl, "Ferramentas & APIs", right_l, Inches(0.97), col_w, Inches(0.28),
    size=Pt(12), bold=True, color=BLUE)
bullet_lines(sl, [
    "Alpaca Markets API — docs.alpaca.markets",
    "Scikit-learn — Pedregosa et al. (2011). JMLR, 12, 2825–2830.",
    "Google Gemini Flash — ai.google.dev",
    "Plotly.js — plotly.com/javascript",
    "Railway Cloud Platform — railway.app",
], right_l, Inches(1.3), col_w, Inches(2.2))

key_msg(sl, "Todos os dados usados são históricos públicos via Alpaca — sem dados proprietários.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Q&A / OBRIGADO
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()

logo_w = Inches(4.5)
logo_h = Inches(1.0)
logo_l = (W - logo_w) / 2
logo_t = Inches(1.1)
box(sl, logo_l, logo_t, logo_w, logo_h, fill_color=RGBColor(0x1a, 0x5c, 0x30), border_color=None)
txt(sl, "ISLA Bot", logo_l, logo_t, logo_w, logo_h,
    size=Pt(40), bold=True, color=RGBColor(0x0d, 0x11, 0x17), align=PP_ALIGN.CENTER)

txt(sl, "Obrigado", Inches(1), Inches(2.3), Inches(11.33), Inches(0.75),
    size=Pt(38), bold=True, color=TEXT, align=PP_ALIGN.CENTER)
txt(sl, "Questões & Discussão", Inches(1), Inches(3.05), Inches(11.33), Inches(0.35),
    size=Pt(14), color=MUTED, align=PP_ALIGN.CENTER)

# final metric cards
summary = [
    ("Retorno Bot",      "−0.50%",   RED,   "6 meses"),
    ("Buy & Hold",       "−10.74%",  RED,   "BTC mesmo período"),
    ("Outperformance",   "+10.25%",  GREEN, "vantagem real"),
    ("Win Rate",         "52.2%",    GREEN, "23 trades"),
]
sw = Inches(2.8)
sg = Inches(0.18)
sh = Inches(1.05)
st = Inches(3.52)
total_sw = 4 * sw + 3 * sg
sl_start = (W - total_sw) / 2
for i, (lbl, val, col, desc) in enumerate(summary):
    card(sl, sl_start + i * (sw + sg), st, sw, sh, lbl, val, desc, col)

txt(sl, "Raphael Malburg · André Neves · Vasco · Beatriz Ferreira\n"
        "Engenharia de Software e IA · Março 2026\n"
        "⚠ Projeto académico · Não é recomendação de investimento · Paper trading apenas",
    Inches(1), Inches(4.75), Inches(11.33), Inches(0.65),
    size=Pt(8.5), color=RGBColor(0x55, 0x5c, 0x64), align=PP_ALIGN.CENTER)


# ── save ────────────────────────────────────────────────────────────────────
out = "isla_bot_presentation.pptx"
prs.save(out)
print(f"Saved: {out}  ({prs.slides.__len__()} slides)")
