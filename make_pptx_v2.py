"""
ISLA Bot — Generate PT and EN presentations (light theme, 13 slides each).
Run:  python make_pptx_v2.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

# ── Palette (GitHub light) ───────────────────────────────────────────────────
BG       = RGBColor(0xFF, 0xFF, 0xFF)
SURFACE  = RGBColor(0xF6, 0xF8, 0xFA)
BORDER   = RGBColor(0xD0, 0xD7, 0xDE)
TEXT     = RGBColor(0x1F, 0x23, 0x28)
MUTED    = RGBColor(0x65, 0x6D, 0x76)
GREEN    = RGBColor(0x1A, 0x7F, 0x37)
DARK_G   = RGBColor(0x0A, 0x50, 0x20)
BLUE     = RGBColor(0x09, 0x69, 0xDA)
RED      = RGBColor(0xCF, 0x22, 0x2E)
AMBER    = RGBColor(0x9A, 0x67, 0x00)
GREEN_BG = RGBColor(0xDF, 0xF0, 0xE5)
BLUE_BG  = RGBColor(0xDB, 0xED, 0xFF)
AMBER_BG = RGBColor(0xFF, 0xF8, 0xC5)
RED_BG   = RGBColor(0xFF, 0xEB, 0xE9)

W = Inches(13.33)
H = Inches(7.5)
FONT = 'Calibri'

# ── Low-level helpers ────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def add_slide(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = BG
    return sl


def fade(slide):
    ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    tr = etree.SubElement(slide.element, f'{{{ns}}}transition')
    tr.set('spd', 'med')
    etree.SubElement(tr, f'{{{ns}}}fade')


def box(sl, l, t, w, h, fill=None, border_color=None, bw=Pt(1)):
    sh = sl.shapes.add_shape(1, l, t, w, h)
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if border_color:
        sh.line.width = bw
        sh.line.color.rgb = border_color
    else:
        sh.line.fill.background()
    return sh


def tb(sl, text, l, t, w, h,
       size=Pt(11), bold=False, italic=False,
       color=TEXT, align=PP_ALIGN.LEFT, wrap=True):
    txb = sl.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size  = size
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name  = FONT
    return txb


def tb_lines(sl, lines, l, t, w, h, default_size=Pt(11), default_color=TEXT):
    """lines: list of str  OR  (str, color, bold, size)"""
    txb = sl.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            item = (item, default_color, False, default_size)
        txt, col, bld, sz = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(2)
        r = p.add_run()
        r.text = txt
        r.font.size  = sz
        r.font.bold  = bld
        r.font.color.rgb = col
        r.font.name  = FONT
    return txb


def bullet_list(sl, items, l, t, w, h, size=Pt(11), color=TEXT, spacing=Pt(4)):
    txb = sl.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, tuple):
            txt, col = item
        else:
            txt, col = item, color
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = spacing
        r = p.add_run()
        r.text = '▸  ' + txt
        r.font.size  = size
        r.font.color.rgb = col
        r.font.name  = FONT
    return txb


# ── Slide-level chrome ───────────────────────────────────────────────────────

def header(sl, title, n=None, total=None, accent=GREEN):
    box(sl, 0, 0, W, Inches(0.07), fill=accent)
    tb(sl, title, Inches(0.45), Inches(0.14), Inches(11.5), Inches(0.55),
       size=Pt(24), bold=True, color=TEXT)
    box(sl, Inches(0.45), Inches(0.72), Inches(12.43), Pt(1.2), fill=BORDER)
    if n and total:
        tb(sl, f'{n} / {total}', Inches(12.0), Inches(0.18), Inches(1.2), Inches(0.3),
           size=Pt(9), color=MUTED, align=PP_ALIGN.RIGHT)
    tb(sl, 'ISLA Bot', Inches(0.45), H - Inches(0.32), Inches(1.5), Inches(0.28),
       size=Pt(8), color=BORDER, bold=True)


def key_bar(sl, text, accent=GREEN):
    bg = {GREEN: GREEN_BG, BLUE: BLUE_BG, AMBER: AMBER_BG, RED: RED_BG}.get(accent, GREEN_BG)
    box(sl, 0, H - Inches(0.5), W, Inches(0.5), fill=bg)
    box(sl, 0, H - Inches(0.5), W, Pt(2.5), fill=accent)
    tb(sl, '\U0001f4a1  ' + text, Inches(0.45), H - Inches(0.44), Inches(12.4), Inches(0.38),
       size=Pt(10.5), bold=True, color=accent)


def card(sl, l, t, w, h, label, value, desc='', vc=GREEN, vs=Pt(30)):
    box(sl, l, t, w, h, fill=SURFACE, border_color=BORDER)
    tb(sl, label.upper(), l + Inches(0.12), t + Inches(0.1),
       w - Inches(0.2), Inches(0.22), size=Pt(7.5), bold=True, color=MUTED)
    tb(sl, value, l + Inches(0.12), t + Inches(0.3),
       w - Inches(0.2), Inches(0.62), size=vs, bold=True, color=vc)
    if desc:
        tb(sl, desc, l + Inches(0.12), t + h - Inches(0.3),
           w - Inches(0.2), Inches(0.26), size=Pt(8.5), color=MUTED)


def callout(sl, text, l, t, w, h, accent=GREEN, size=Pt(10.5)):
    bg = {GREEN: GREEN_BG, BLUE: BLUE_BG, AMBER: AMBER_BG, RED: RED_BG}.get(accent, GREEN_BG)
    box(sl, l, t, w, h, fill=bg)
    box(sl, l, t, Pt(5), h, fill=accent)
    tb(sl, text, l + Inches(0.15), t + Inches(0.1),
       w - Inches(0.22), h - Inches(0.16), size=size, color=TEXT, wrap=True)


def section_label(sl, text, l, t, w, accent=GREEN):
    box(sl, l, t, Pt(4), Inches(0.24), fill=accent)
    tb(sl, text.upper(), l + Inches(0.12), t, w - Inches(0.15), Inches(0.24),
       size=Pt(8.5), bold=True, color=accent)


def htable(sl, rows, l, t, w, row_h=Inches(0.34),
           header_fill=SURFACE, header_color=MUTED, col_widths=None):
    """rows: list of list of (text, color, bold).  First row = header."""
    n_cols = len(rows[0])
    if col_widths is None:
        col_widths = [w / n_cols] * n_cols
    cur_t = t
    for ri, row in enumerate(rows):
        is_hdr = ri == 0
        bg = header_fill if is_hdr else (SURFACE if ri % 2 == 1 else BG)
        box(sl, l, cur_t, w, row_h, fill=bg, border_color=BORDER)
        cur_l = l
        for ci, cell in enumerate(row):
            if isinstance(cell, str):
                cell = (cell, MUTED if is_hdr else TEXT, is_hdr)
            ct, cc, cb = cell
            cw = col_widths[ci]
            tb(sl, ct, cur_l + Inches(0.1), cur_t + Inches(0.06),
               cw - Inches(0.12), row_h - Inches(0.1),
               size=Pt(8.5 if is_hdr else 10), color=cc, bold=cb)
            cur_l += cw
        cur_t += row_h
    return cur_t  # bottom y


# ════════════════════════════════════════════════════════════════════════════
# CONTENT — both languages
# ════════════════════════════════════════════════════════════════════════════

CONTENT = {
    'PT': {
        'cover': {
            'title':    'ISLA Bot',
            'subtitle': 'Trading Algorítmico com Machine Learning',
            'tags':     'BTC/USD  ·  Paper Trading  ·  Random Forest  ·  Gemini Flash',
            'group_label': 'Grupo',
            'members':  'Raphael Malburg   André Neves   Vasco   Beatriz Ferreira',
            'course':   'Engenharia de Software e IA  ·  Março 2026',
            'disclaimer': '⚠️  Projeto académico · Não é recomendação de investimento · Paper trading apenas',
        },
        'problem': {
            'title': 'Motivação e Problema',
            'stats': [
                ('24/7', RED,   'Mercados cripto nunca fecham—monitorização humana contínua é impossível'),
                ('≤70%', AMBER, 'dos traders individuais perdem\ndinheiro por decisões emocionais'),
                ('ML',   GREEN, 'Dados históricos ricos permitem\nmodelar padrões sem emoção'),
            ],
            'question_label': 'Pergunta de investigação',
            'question': (
                'É possível construir um sistema algorítmico que identifique oportunidades em BTC/USD '
                'e supere Buy & Hold, usando 17 indicadores técnicos e Machine Learning, '
                'com gestão de risco automática e sentimento de mercado como gate?'
            ),
            'bar': 'A vantagem do bot não é ser mais inteligente—é ser consistente e disciplinado.',
        },
        'solution': {
            'title': 'A Solução — ISLA Bot',
            'bullets': [
                'Bot de trading algorítmico para BTC/USD em paper trading',
                'Candles de 1H analisados com 17 indicadores técnicos',
                'Random Forest prevê se o Take-Profit é atingido antes do Stop-Loss',
                'Gate de sentimento via Gemini Flash (OpenRouter) bloqueia entradas bearish',
                'Exit watchdog a cada 10 s: SL (1×ATR), TP (2×ATR), máx. 12 candles',
                'Dashboard web em tempo real com gráficos Plotly.js (polling 15 s)',
            ],
            'specs': [
                ('Exchange',        'Alpaca Markets', 'Paper trading, sem risco real'),
                ('Stack',           'Python · Flask · Scikit-learn · SQLite', ''),
                ('LLM / Sentiment', 'Gemini Flash via OpenRouter', ''),
                ('Ativo / Timeframe', 'BTC/USD  ·  1H candles', ''),
            ],
            'bar': 'Pipeline end-to-end: dados → ML → sentimento → risco → ordem na exchange.',
        },
        'architecture': {
            'title': 'Arquitetura do Sistema',
            'layers': [
                ('Fontes de Dados',
                 [('Alpaca API', 'Candles OHLCV · Ordens · Conta', BLUE),
                  ('OpenRouter / Gemini Flash', 'Sentimento de mercado LLM', AMBER)]),
                ('Pipeline Core',
                 [('src/data.py', 'Fetch + 17 indicadores', MUTED),
                  ('src/model.py', 'Random Forest · predict_proba', BLUE),
                  ('src/trading.py', 'Orquestração + Sentimento', MUTED),
                  ('src/strategy.py', 'Sizing · SL/TP · Circuit Breaker', MUTED),
                  ('src/broker.py', 'Alpaca: ordens + posições', GREEN)]),
                ('Orquestração',
                 [('src/app.py (Flask)', 'bot_loop (1H) · exit_watchdog (10 s) · REST API · Threads', BLUE),
                  ('SQLite', 'runs · trades · equity history', MUTED)]),
                ('Interface',
                 [('Dashboard Web', 'Jinja2 · Plotly.js · Polling 15 s · Kill switch', GREEN)]),
            ],
            'bar': 'Dois threads em background: bot_loop (1H) e exit_watchdog (10 s)—o servidor Flask nunca bloqueia.',
        },
        'ml_model': {
            'title': 'Modelo de Machine Learning',
            'obj_label': 'Objetivo do modelo',
            'obj': '"O Take-Profit (2×ATR) será atingido antes do Stop-Loss (1×ATR) nos próximos 12 candles?"',
            'algo_label': 'Algoritmo',
            'algo': [
                'Random Forest · 200 árvores · profundidade máx. 8',
                'StandardScaler em todos os 17 features',
                'Walk-forward validation — TimeSeriesSplit 5 folds',
                'Sem data leakage temporal',
                'Treino: 6.967 amostras  ·  Teste: 1.742 amostras',
            ],
            'feat_label': 'Top features (importância)',
            'features': [
                ('#', 'Feature', 'Peso'),
                ('1', 'OBV (On-Balance Volume)', '11.9 %'),
                ('2', 'ATR-14 (volatilidade)', '10.0 %'),
                ('3', 'EMA-20', '9.6 %'),
                ('4', 'EMA-50', '8.9 %'),
                ('5', 'BB High', '8.6 %'),
                ('6', 'BB Low', '7.9 %'),
                ('7', 'BB Width', '7.8 %'),
            ],
            'insight': 'Volume (OBV) e volatilidade (ATR, BB) dominam—onde há volume e expansão há movimento.',
            'bar': 'Walk-forward garante que cada previsão usa apenas dados do passado—como numa operação real.',
        },
        'pipeline': {
            'title': 'Pipeline de Decisão — A cada candle de 1H',
            'col_a_label': 'Análise (passos 1–5)',
            'col_b_label': 'Execução (passos 6–8)',
            'steps_a': [
                ('1', 'Fetch Market Data', '90 dias · 1H · ~2.160 barras BTC/USD via Alpaca'),
                ('2', 'Feature Engineering', 'RSI-14, MACD, EMA-20/50, BB, ATR-14, OBV + 11 mais (17 total)'),
                ('3', 'ML Prediction', 'P(TP antes SL em 12H) ≥ 55% → LONG · senão → FLAT'),
                ('4', 'Drift Detection', 'Confiança média das últimas 30 runs < 50% → aviso de drift'),
                ('5', 'Sentiment Gate + Circuit Breaker', 'Score Gemini ≤ −0.5 bloqueia · perda diária > 10% para trading'),
            ],
            'steps_b': [
                ('6', 'Position Sizing', 'equity × 5% ÷ stop_pct  ·  stop_pct = 1×ATR / preço'),
                ('7', 'Market BUY + Disaster Stop', 'Ordem de mercado + stop-limit a 3×ATR (safety net de crash)'),
                ('8', 'Exit Watchdog — a cada 10 s', 'SL (1×ATR) · TP (2×ATR) · max 12H · sinal flip → fechar'),
            ],
            'watchdog_note': (
                'O watchdog corre a cada 10 s em background e não depende da exchange—'
                'elimina ordens órfãs e permite saídas baseadas em qualquer condição.'
            ),
            'bar': 'O bot entra pouco e sai rápido—selectividade é a principal fonte de vantagem.',
        },
        'risk': {
            'title': 'Gestão de Risco',
            'items': [
                ('\U0001f7e2\U0001f534', 'SL (1×ATR) e TP (2×ATR)',
                 'Rácio risco/retorno de 2:1. Geridos pelo exit watchdog a cada 10 s—não por ordens na exchange. Uma posição vencedora compensa duas perdedoras.'),
                ('⚡', 'Circuit Breaker + Max Hold',
                 'Bloqueia novos trades se perda diária ultrapassar 10% do capital. Fecha posição automaticamente após 12 candles (12H) mesmo sem SL/TP.'),
                ('⚠️', 'Disaster Stop (3×ATR)',
                 'Stop-limit colocado na Alpaca a 3×ATR abaixo da entrada. Dispara apenas se o servidor cair. Protege contra crash catastrófico sem interferir com saídas normais.'),
                ('\U0001f4ca', 'Position Sizing (5% fixo)',
                 'Risco fixo de 5% do capital por trade, ajustado pela volatilidade (ATR). Testámos Kelly Criterion—não melhora com edge negativo por trade.'),
            ],
            'bar': 'Duas camadas de proteção: app-side (lógica real) + exchange-side (safety net)—resiliente a falhas.',
        },
        'metrics': {
            'title': 'Métricas do Modelo (holdout 20%)',
            'cards': [
                ('Accuracy',  '67.9%', 'previsões corretas',     GREEN),
                ('Precision', '23.6%', 'taxa de vitória real',   AMBER),
                ('Recall',    '10.9%', 'oportunidades captadas', AMBER),
                ('F1 Score',  '0.149', 'P × R balanceados',     MUTED),
            ],
            'interp_label': 'O que significam estes números',
            'interp': [
                'Accuracy de 68% deve-se ao modelo acertar a maioria dos FLATs (classe dominante)',
                'Precision de 23.6%: o bot ganha ~1 em cada 4 trades que abre',
                'Break-even para RR 2:1 é 33.3%—estamos abaixo do threshold de lucro por trade',
                'Recall baixo = entradas raras = bot altamente selectivo',
            ],
            'train_label': 'Dados',
            'train_rows': [
                ('Amostras treino', '6.967'),
                ('Amostras teste',  '1.742'),
                ('Features',        '17'),
                ('Árvores RF',      '200 · max_depth 8'),
                ('Validação',       'Walk-forward 5 folds'),
            ],
            'bar': 'Precision abaixo de 33.3% = edge negativo por trade—o ganho vem da selectividade, não da precisão.',
            'bar_accent': AMBER,
        },
        'results': {
            'title': 'Resultados — Backtest Walk-Forward 6 Meses',
            'hero_label': 'Outperformance vs Buy & Hold',
            'hero_value': '+10.25%',
            'hero_sub':   'estratégia −0.50%  ·  buy & hold −10.74%',
            'cards': [
                ('Win Rate', '52.2%', '12 de 23 trades', GREEN),
                ('Trades',   '23',    '6 meses  ·  1H candles', TEXT),
            ],
            'table': [
                ['Métrica', 'Valor'],
                [('Max Drawdown',  TEXT, False), ('−12.34%', RED,   True)],
                [('Profit Factor', TEXT, False), ('1.05',        TEXT,  False)],
                [('Avg Win',       TEXT, False), ('+$1.23',      GREEN, True)],
                [('Avg Loss',      TEXT, False), ('−$1.28', RED,   False)],
                [('Sharpe Ratio',  TEXT, False), ('−0.05',  AMBER, False)],
            ],
            'context': (
                'Período fortemente bearish para BTC (−10.7%). '
                'O bot preservou capital ao ficar maioritariamente FLAT—'
                'a vantagem vem de não perder, não de ganhar muito.'
            ),
            'bar': 'Em mercado em baixa, não perder já é ganhar—o bot superou o benchmark em +10.25%.',
        },
        'kelly': {
            'title': 'Análise — Kelly Criterion vs Fixed 5%',
            'table': [
                ['Estratégia', 'Retorno', 'Sharpe', 'Max Drawdown'],
                [('Fixed 5% (atual)', GREEN, True),  ('+4.72%', GREEN, True), ('0.17', TEXT, False), ('22.1%', RED, False)],
                [('Full Kelly',       TEXT,  False), ('+3.11%', GREEN, False),('0.19', TEXT, False), ('10.5%', GREEN, False)],
                [('Half Kelly',       TEXT,  False), ('+1.52%', AMBER, False),('0.16', TEXT, False), ('5.7%',  GREEN, False)],
                [('Quarter Kelly',    TEXT,  False), ('+0.79%', AMBER, False),('0.16', TEXT, False), ('2.9%',  GREEN, False)],
            ],
            'formula_label': 'Fórmula de Kelly',
            'formula': 'f· = (p × b − q) / b',
            'calc':    'Com p=0.236, b=2  ⇒  f· = −14.6% (não apostar)',
            'why':     (
                'O modelo prevê confiança de ~55–65% mas a precisão real é 23.6%. '
                'Kelly penaliza edges negativos de forma agressiva—abaixo de 33.3% de precision, '
                'Fixed Fraction supera sempre Kelly.'
            ),
            'when': 'Kelly funciona quando precision ≥ 33.3% (break-even para RR 2:1). Meta para próxima versão.',
            'bar': 'Kelly só é válido com edge positivo—primeiro melhorar o modelo, depois aplicar Kelly.',
            'bar_accent': AMBER,
        },
        'dashboard': {
            'title': 'Dashboard Web em Tempo Real',
            'features': [
                'Preço BTC e sinal do bot actualizados a cada 15 s sem page refresh',
                'Gráfico de velas interactivo com zoom, pan e scroll (estilo TradingView)',
                'Linhas de SL e TP sobrepostas no gráfico com etiquetas de preço ao vivo',
                'RSI-14 e MACD sincronizados com o gráfico principal',
                'Kill switch para fechar todas as posições imediatamente',
                'Secção de backtest: equity curve + log dos 23 trades',
                'Métricas ML + gráfico de feature importance interactivo',
            ],
            'intervals': [
                ('15 s', 'preço + posição + linhas SL/TP no gráfico'),
                ('30 s', 'sinal ML + trades + métricas completas'),
            ],
            'api': [
                ('GET',  '/api/live_stats',     GREEN),
                ('GET',  '/api/dashboard_data', GREEN),
                ('GET',  '/api/showcase',        GREEN),
                ('POST', '/api/kill_switch',     RED),
            ],
            'note': 'Sem WebSockets—polling assíncrono com fetch(). Plotly.js actualiza gráficos in-place, preservando o zoom do utilizador.',
            'bar': 'Uma única página mostra tudo—o dashboard é a cabine de pilotagem do bot.',
        },
        'team': {
            'title': 'Metodologia & Equipa',
            'process_label': 'Processo de Desenvolvimento',
            'process': [
                'Abordagem iterativa: ciclos semanais de build → test → backtest → refine',
                'Controlo de versões via Git—commits atómicos por feature',
                'Testes de integração com dados reais da Alpaca (paper trading)',
                'Walk-forward backtest como critério de aceitação de cada modelo',
                'Deploy contínuo em Railway—dashboard sempre acessível ao grupo',
            ],
            'rule': 'Nenhuma feature considerada "pronta" sem backtest validar métricas iguais ou superiores à iteração anterior.',
            'team_label': 'Contribuições',
            'members': [
                ('Raphael Malburg',  'Arquitetura · Flask API · Dashboard · Deploy Railway · Integração Alpaca'),
                ('André Neves',      'Pipeline ML · Feature engineering · Backtesting · Análise de métricas'),
                ('Vasco',            'Gestão de risco · Exit watchdog · Análise Kelly Criterion · Position sizing'),
                ('Beatriz Ferreira', 'Integração Gemini/sentimento · Testes integração · Documentação'),
            ],
            'bar': 'Cada módulo foi desenvolvido e testado de forma independente—falhas identificadas cedo.',
        },
        'conclusions': {
            'title': 'Conclusões & Próximos Passos',
            'wins_label': 'O que conseguimos',
            'wins': [
                '+10.25% de outperformance vs Buy & Hold num período bearish',
                'Pipeline completo end-to-end com 8 camadas de decisão',
                'Exit management robusto: watchdog 10 s + disaster stop 3×ATR',
                'Dashboard profissional com actualização em tempo real',
                'Deploy operacional 24/7 em Railway',
            ],
            'limits_label': 'Limitações actuais',
            'limits': [
                'Precision 23.6%—abaixo do break-even de 33.3% (RR 2:1)',
                'Apenas 23 trades em 6 meses—frequência de sinal baixa',
                'Kelly Criterion não aplicável sem edge positivo por trade',
                'Paper trading—slippage real pode ser maior',
            ],
            'next_label': 'Próximos passos',
            'next': [
                'Substituir RF por XGBoost / LightGBM—melhor para séries temporais financeiras',
                'Aumentar threshold de confiança para ≥ 65%—menos trades, precision maior',
                'Adicionar funding rates e open interest como features',
                'Trailing stop para proteger lucros em tendências prolongadas',
                'Alertas em tempo real via Telegram / Discord',
            ],
            'summary': (
                'O bot demonstra que selectividade inteligente—ficar FLAT na maioria das barras—'
                'pode superar Buy & Hold em mercados em baixa. '
                'O valor está na ausência de maus trades, não na qualidade de cada trade individual.'
            ),
            'bar': 'Um passo de cada vez: primeiro precision > 33%, depois Kelly, depois capital real.',
        },
        'closing': {
            'title':      'ISLA Bot',
            'thanks':     'Obrigado',
            'qa':         'Questões & Discussão',
            'outperform': ('+10.25%', 'Outperformance vs B&H'),
            'strategy':   ('−0.50%',  'Retorno do bot'),
            'bh':         ('−10.74%', 'Buy & Hold BTC'),
            'winrate':    ('52.2%',   'Win Rate'),
            'team':       'Raphael Malburg  ·  André Neves  ·  Vasco  ·  Beatriz Ferreira',
            'course':     'Engenharia de Software e IA  ·  Março 2026',
            'disclaimer': '⚠️  Projeto académico · Não é recomendação de investimento · Paper trading apenas',
        },
    },
    'EN': {
        'cover': {
            'title':    'ISLA Bot',
            'subtitle': 'Algorithmic Trading with Machine Learning',
            'tags':     'BTC/USD  ·  Paper Trading  ·  Random Forest  ·  Gemini Flash',
            'group_label': 'Team',
            'members':  'Raphael Malburg   André Neves   Vasco   Beatriz Ferreira',
            'course':   'Software Engineering & AI  ·  March 2026',
            'disclaimer': '⚠️  Academic project · Not investment advice · Paper trading only',
        },
        'problem': {
            'title': 'Motivation & Problem',
            'stats': [
                ('24/7', RED,   'Crypto markets never close—continuous human monitoring is impossible'),
                ('≤70%', AMBER, 'of individual traders lose money\ndue to emotional and inconsistent decisions'),
                ('ML',   GREEN, 'Rich historical data enables\npattern modelling without emotion'),
            ],
            'question_label': 'Research question',
            'question': (
                'Can we build an algorithmic system that identifies BTC/USD opportunities '
                'and outperforms Buy & Hold, using 17 technical indicators and Machine Learning, '
                'with automatic risk management and a market sentiment gate?'
            ),
            'bar': 'The bot\'s edge is not being smarter—it\'s being consistent and disciplined.',
        },
        'solution': {
            'title': 'The Solution — ISLA Bot',
            'bullets': [
                'Algorithmic trading bot for BTC/USD via paper trading',
                '1H candles analysed with 17 technical indicators',
                'Random Forest predicts whether Take-Profit is hit before Stop-Loss',
                'Sentiment gate via Gemini Flash (OpenRouter) blocks bearish entries',
                'Exit watchdog every 10 s: SL (1×ATR), TP (2×ATR), max 12 candles',
                'Real-time web dashboard with Plotly.js charts (15 s polling)',
            ],
            'specs': [
                ('Exchange',         'Alpaca Markets', 'Paper trading, no real funds'),
                ('Stack',            'Python · Flask · Scikit-learn · SQLite', ''),
                ('LLM / Sentiment',  'Gemini Flash via OpenRouter', ''),
                ('Asset / Timeframe','BTC/USD  ·  1H candles', ''),
            ],
            'bar': 'End-to-end pipeline: data → ML → sentiment → risk → order on exchange.',
        },
        'architecture': {
            'title': 'System Architecture',
            'layers': [
                ('Data Sources',
                 [('Alpaca API', 'OHLCV Candles · Orders · Account', BLUE),
                  ('OpenRouter / Gemini Flash', 'LLM market sentiment', AMBER)]),
                ('Core Pipeline',
                 [('src/data.py', 'Fetch + 17 indicators', MUTED),
                  ('src/model.py', 'Random Forest · predict_proba', BLUE),
                  ('src/trading.py', 'Orchestration + Sentiment', MUTED),
                  ('src/strategy.py', 'Sizing · SL/TP · Circuit Breaker', MUTED),
                  ('src/broker.py', 'Alpaca: orders + positions', GREEN)]),
                ('Orchestration',
                 [('src/app.py (Flask)', 'bot_loop (1H) · exit_watchdog (10 s) · REST API · Threads', BLUE),
                  ('SQLite', 'runs · trades · equity history', MUTED)]),
                ('Interface',
                 [('Web Dashboard', 'Jinja2 · Plotly.js · 15 s polling · Kill switch', GREEN)]),
            ],
            'bar': 'Two background threads: bot_loop (1H) and exit_watchdog (10 s)—Flask server never blocks.',
        },
        'ml_model': {
            'title': 'Machine Learning Model',
            'obj_label': 'Model objective',
            'obj': '"Will Take-Profit (2×ATR) be hit before Stop-Loss (1×ATR) in the next 12 candles?"',
            'algo_label': 'Algorithm',
            'algo': [
                'Random Forest · 200 trees · max depth 8',
                'StandardScaler on all 17 features',
                'Walk-forward validation — TimeSeriesSplit 5 folds',
                'No temporal data leakage',
                'Train: 6,967 samples  ·  Test: 1,742 samples',
            ],
            'feat_label': 'Top features (importance)',
            'features': [
                ('#', 'Feature', 'Weight'),
                ('1', 'OBV (On-Balance Volume)', '11.9 %'),
                ('2', 'ATR-14 (volatility)', '10.0 %'),
                ('3', 'EMA-20', '9.6 %'),
                ('4', 'EMA-50', '8.9 %'),
                ('5', 'BB High', '8.6 %'),
                ('6', 'BB Low', '7.9 %'),
                ('7', 'BB Width', '7.8 %'),
            ],
            'insight': 'Volume (OBV) and volatility (ATR, BB) dominate—where there is volume and expansion there is movement.',
            'bar': 'Walk-forward ensures each prediction uses only past data—as in a real live operation.',
        },
        'pipeline': {
            'title': 'Decision Pipeline — Every 1H candle',
            'col_a_label': 'Analysis (steps 1–5)',
            'col_b_label': 'Execution (steps 6–8)',
            'steps_a': [
                ('1', 'Fetch Market Data', '90 days · 1H · ~2,160 BTC/USD bars via Alpaca'),
                ('2', 'Feature Engineering', 'RSI-14, MACD, EMA-20/50, BB, ATR-14, OBV + 11 more (17 total)'),
                ('3', 'ML Prediction', 'P(TP before SL in 12H) ≥ 55% → LONG · else → FLAT'),
                ('4', 'Drift Detection', 'Avg confidence of last 30 runs < 50% → drift warning'),
                ('5', 'Sentiment Gate + Circuit Breaker', 'Gemini score ≤ −0.5 blocks entry · daily loss > 10% halts trading'),
            ],
            'steps_b': [
                ('6', 'Position Sizing', 'equity × 5% ÷ stop_pct  ·  stop_pct = 1×ATR / price'),
                ('7', 'Market BUY + Disaster Stop', 'Market order + stop-limit at 3×ATR (crash safety net)'),
                ('8', 'Exit Watchdog — every 10 s', 'SL (1×ATR) · TP (2×ATR) · max 12H · signal flip → close'),
            ],
            'watchdog_note': (
                'The watchdog runs every 10 s in background and is independent of the exchange—'
                'eliminates orphaned orders and allows exits based on any condition.'
            ),
            'bar': 'The bot enters rarely and exits fast—selectivity is the main source of edge.',
        },
        'risk': {
            'title': 'Risk Management',
            'items': [
                ('\U0001f7e2\U0001f534', 'SL (1×ATR) and TP (2×ATR)',
                 'Risk/reward ratio of 2:1. Managed by exit watchdog every 10 s—not by exchange orders. One winning trade covers two losers.'),
                ('⚡', 'Circuit Breaker + Max Hold',
                 'Blocks new trades if daily loss exceeds 10% of capital. Automatically closes position after 12 candles (12H) even without SL/TP.'),
                ('⚠️', 'Disaster Stop (3×ATR)',
                 'Stop-limit placed on Alpaca at 3×ATR below entry. Fires only if the server crashes. Protects against catastrophic loss without interfering with normal exits.'),
                ('\U0001f4ca', 'Position Sizing (fixed 5%)',
                 'Fixed risk of 5% of capital per trade, adjusted by volatility (ATR). Kelly Criterion was tested—it does not improve results with a negative per-trade edge.'),
            ],
            'bar': 'Two protection layers: app-side (real logic) + exchange-side (safety net)—resilient to failures.',
        },
        'metrics': {
            'title': 'Model Metrics (20% holdout)',
            'cards': [
                ('Accuracy',  '67.9%', 'correct predictions',      GREEN),
                ('Precision', '23.6%', 'real win rate',             AMBER),
                ('Recall',    '10.9%', 'opportunities captured',    AMBER),
                ('F1 Score',  '0.149', 'P × R balanced',           MUTED),
            ],
            'interp_label': 'What these numbers mean',
            'interp': [
                '68% accuracy comes from correctly predicting the dominant FLAT class',
                'Precision 23.6%: bot wins ~1 out of every 4 trades it opens',
                'Break-even for RR 2:1 is 33.3%—we are below the per-trade profit threshold',
                'Low recall = rare entries = highly selective bot',
            ],
            'train_label': 'Data',
            'train_rows': [
                ('Train samples', '6,967'),
                ('Test samples',  '1,742'),
                ('Features',      '17'),
                ('RF trees',      '200 · max_depth 8'),
                ('Validation',    'Walk-forward 5 folds'),
            ],
            'bar': 'Precision below 33.3% = negative per-trade edge—the gain comes from selectivity, not precision.',
            'bar_accent': AMBER,
        },
        'results': {
            'title': 'Results — 6-Month Walk-Forward Backtest',
            'hero_label': 'Outperformance vs Buy & Hold',
            'hero_value': '+10.25%',
            'hero_sub':   'strategy −0.50%  ·  buy & hold −10.74%',
            'cards': [
                ('Win Rate', '52.2%', '12 of 23 trades', GREEN),
                ('Trades',   '23',    '6 months  ·  1H candles', TEXT),
            ],
            'table': [
                ['Metric', 'Value'],
                [('Max Drawdown',  TEXT, False), ('−12.34%', RED,   True)],
                [('Profit Factor', TEXT, False), ('1.05',        TEXT,  False)],
                [('Avg Win',       TEXT, False), ('+$1.23',      GREEN, True)],
                [('Avg Loss',      TEXT, False), ('−$1.28', RED,   False)],
                [('Sharpe Ratio',  TEXT, False), ('−0.05',  AMBER, False)],
            ],
            'context': (
                'Strongly bearish period for BTC (−10.7%). '
                'The bot preserved capital by staying mostly FLAT—'
                'the advantage comes from NOT losing, not from big gains.'
            ),
            'bar': 'In a bear market, not losing is already winning—the bot outperformed the benchmark by +10.25%.',
        },
        'kelly': {
            'title': 'Analysis — Kelly Criterion vs Fixed 5%',
            'table': [
                ['Strategy', 'Return', 'Sharpe', 'Max Drawdown'],
                [('Fixed 5% (current)', GREEN, True),  ('+4.72%', GREEN, True), ('0.17', TEXT, False), ('22.1%', RED, False)],
                [('Full Kelly',         TEXT,  False), ('+3.11%', GREEN, False),('0.19', TEXT, False), ('10.5%', GREEN, False)],
                [('Half Kelly',         TEXT,  False), ('+1.52%', AMBER, False),('0.16', TEXT, False), ('5.7%',  GREEN, False)],
                [('Quarter Kelly',      TEXT,  False), ('+0.79%', AMBER, False),('0.16', TEXT, False), ('2.9%',  GREEN, False)],
            ],
            'formula_label': 'Kelly formula',
            'formula': 'f· = (p × b − q) / b',
            'calc':    'With p=0.236, b=2  ⇒  f· = −14.6% (do not bet)',
            'why':     (
                'The model predicts confidence of ~55–65% but real precision is 23.6%. '
                'Kelly penalises negative edges aggressively—below 33.3% precision, '
                'Fixed Fraction always outperforms Kelly.'
            ),
            'when': 'Kelly works when precision ≥ 33.3% (break-even for RR 2:1). Target for next version.',
            'bar': 'Kelly is only valid with a positive per-trade edge—improve the model first, then apply Kelly.',
            'bar_accent': AMBER,
        },
        'dashboard': {
            'title': 'Real-Time Web Dashboard',
            'features': [
                'BTC price and bot signal updated every 15 s without page refresh',
                'Interactive candlestick chart with zoom, pan and scroll (TradingView-style)',
                'SL and TP lines overlaid on the chart with live price labels',
                'RSI-14 and MACD synchronised with the main chart',
                'Kill switch to close all positions immediately',
                'Backtest section: 6-month equity curve + log of all 23 trades',
                'ML metrics + interactive feature importance chart',
            ],
            'intervals': [
                ('15 s', 'price + position + SL/TP lines on chart'),
                ('30 s', 'ML signal + trades + full metrics'),
            ],
            'api': [
                ('GET',  '/api/live_stats',     GREEN),
                ('GET',  '/api/dashboard_data', GREEN),
                ('GET',  '/api/showcase',        GREEN),
                ('POST', '/api/kill_switch',     RED),
            ],
            'note': 'No WebSockets—async polling with fetch(). Plotly.js updates charts in-place, preserving user zoom.',
            'bar': 'One page shows everything—the dashboard is the bot\'s cockpit.',
        },
        'team': {
            'title': 'Methodology & Team',
            'process_label': 'Development Process',
            'process': [
                'Iterative approach: weekly build → test → backtest → refine cycles',
                'Version control via Git—atomic commits per feature',
                'Integration tests with real Alpaca data (paper trading)',
                'Walk-forward backtest as acceptance criterion for each model',
                'Continuous deployment on Railway—dashboard always accessible to the team',
            ],
            'rule': 'No feature considered done without a backtest validating equal or better metrics than the previous iteration.',
            'team_label': 'Contributions',
            'members': [
                ('Raphael Malburg',  'Architecture · Flask API · Dashboard · Railway Deploy · Alpaca Integration'),
                ('André Neves',      'ML Pipeline · Feature engineering · Backtesting · Metrics analysis'),
                ('Vasco',            'Risk management · Exit watchdog · Kelly Criterion analysis · Position sizing'),
                ('Beatriz Ferreira', 'Gemini/sentiment integration · Integration tests · Documentation'),
            ],
            'bar': 'Each module was developed and tested independently—failures identified early.',
        },
        'conclusions': {
            'title': 'Conclusions & Next Steps',
            'wins_label': 'What we achieved',
            'wins': [
                '+10.25% outperformance vs Buy & Hold in a bearish period',
                'Complete end-to-end pipeline with 8 decision layers',
                'Robust exit management: 10 s watchdog + 3×ATR disaster stop',
                'Professional dashboard with real-time updates',
                '24/7 operational deployment on Railway',
            ],
            'limits_label': 'Current limitations',
            'limits': [
                'Precision 23.6%—below break-even threshold of 33.3% (RR 2:1)',
                'Only 23 trades in 6 months—low signal frequency',
                'Kelly Criterion not applicable without positive per-trade edge',
                'Paper trading—real slippage may be higher',
            ],
            'next_label': 'Next steps',
            'next': [
                'Replace RF with XGBoost / LightGBM—better for financial time series',
                'Raise confidence threshold to ≥ 65%—fewer trades, higher precision',
                'Add funding rates and open interest as features',
                'Trailing stop to lock in profits during extended trends',
                'Real-time alerts via Telegram / Discord',
            ],
            'summary': (
                'The bot demonstrates that intelligent selectivity—staying FLAT most of the time—'
                'can outperform Buy & Hold in bear markets. '
                'The value lies in the absence of bad trades, not in the quality of each individual trade.'
            ),
            'bar': 'One step at a time: precision > 33% first, then Kelly, then real capital.',
        },
        'closing': {
            'title':      'ISLA Bot',
            'thanks':     'Thank you',
            'qa':         'Questions & Discussion',
            'outperform': ('+10.25%', 'Outperformance vs B&H'),
            'strategy':   ('−0.50%',  'Strategy return'),
            'bh':         ('−10.74%', 'Buy & Hold BTC'),
            'winrate':    ('52.2%',   'Win Rate'),
            'team':       'Raphael Malburg  ·  André Neves  ·  Vasco  ·  Beatriz Ferreira',
            'course':     'Software Engineering & AI  ·  March 2026',
            'disclaimer': '⚠️  Academic project · Not investment advice · Paper trading only',
        },
    },
}

TOTAL = 13

# ════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ════════════════════════════════════════════════════════════════════════════

def slide_cover(prs, c):
    sl = add_slide(prs)
    fade(sl)
    # Left green panel
    panel_w = Inches(5.1)
    box(sl, 0, 0, panel_w, H, fill=DARK_G)
    # decorative circles
    box(sl, Inches(-0.8), Inches(-0.8), Inches(2.5), Inches(2.5),
        fill=RGBColor(0x0E, 0x60, 0x28))
    box(sl, Inches(3.5), Inches(5.8), Inches(2.5), Inches(2.5),
        fill=RGBColor(0x0E, 0x60, 0x28))
    # Logo text
    tb(sl, c['title'], Inches(0.3), Inches(2.0), panel_w - Inches(0.5), Inches(1.4),
       size=Pt(56), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    # Thin white underline
    box(sl, Inches(0.6), Inches(3.45), panel_w - Inches(1.1), Pt(2),
        fill=RGBColor(0x7F, 0xD8, 0x9F))
    tb(sl, c['subtitle'], Inches(0.2), Inches(3.6), panel_w - Inches(0.3), Inches(0.55),
       size=Pt(12.5), color=RGBColor(0xAD, 0xE8, 0xC3), align=PP_ALIGN.CENTER, wrap=True)
    tb(sl, c['tags'], Inches(0.2), Inches(4.22), panel_w - Inches(0.3), Inches(0.35),
       size=Pt(9), color=RGBColor(0x7F, 0xD8, 0x9F), align=PP_ALIGN.CENTER, wrap=True)
    # Right side
    rl = panel_w + Inches(0.6)
    rw = W - rl - Inches(0.5)
    tb(sl, c['group_label'].upper(), rl, Inches(1.6), rw, Inches(0.25),
       size=Pt(9), bold=True, color=MUTED)
    box(sl, rl, Inches(1.88), rw, Pt(1.5), fill=BORDER)
    tb(sl, c['members'], rl, Inches(1.98), rw, Inches(0.6),
       size=Pt(14), bold=True, color=TEXT, wrap=True)
    box(sl, rl, Inches(2.7), Inches(3.8), Pt(1.5), fill=GREEN)
    # Group box
    grp_t = Inches(2.85)
    grp_h = Inches(1.25)
    box(sl, rl, grp_t, rw, grp_h, fill=SURFACE, border_color=BORDER)
    tb(sl, c['course'], rl + Inches(0.15), grp_t + Inches(0.15), rw - Inches(0.25), Inches(0.32),
       size=Pt(11.5), color=TEXT)
    tb(sl, c['disclaimer'], rl + Inches(0.15), grp_t + Inches(0.56), rw - Inches(0.25), Inches(0.55),
       size=Pt(9.5), color=MUTED, wrap=True)
    return sl


def slide_problem(prs, c, n):
    sl = add_slide(prs)
    fade(sl)
    header(sl, c['title'], n=n, total=TOTAL)
    sw = Inches(4.0)
    sg = Inches(0.2)
    total_w = 3 * sw + 2 * sg
    sl_start = (W - total_w) / 2
    st = Inches(1.0)
    sh = Inches(2.3)
    for i, (num, col, desc) in enumerate(c['stats']):
        l = sl_start + i * (sw + sg)
        box(sl, l, st, sw, sh, fill=SURFACE, border_color=BORDER)
        # colored top accent
        box(sl, l, st, sw, Inches(0.06), fill=col)
        tb(sl, num, l, st + Inches(0.15), sw, Inches(1.0),
           size=Pt(52), bold=True, color=col, align=PP_ALIGN.CENTER)
        tb(sl, desc, l + Inches(0.15), st + Inches(1.2), sw - Inches(0.25), Inches(0.95),
           size=Pt(10.5), color=MUTED, align=PP_ALIGN.CENTER, wrap=True)
    # Research question
    rq_t = Inches(3.5)
    tb(sl, c['question_label'].upper(), Inches(0.45), rq_t - Inches(0.28), Inches(5), Inches(0.24),
       size=Pt(8.5), bold=True, color=GREEN)
    callout(sl, c['question'], Inches(0.45), rq_t, Inches(12.43), Inches(1.1), accent=GREEN)
    key_bar(sl, c['bar'])
    return sl


def slide_solution(prs, c, n):
    sl = add_slide(prs)
    fade(sl)
    header(sl, c['title'], n=n, total=TOTAL)
    col_w = Inches(6.7)
    bullet_list(sl, c['bullets'], Inches(0.45), Inches(0.92), col_w, Inches(3.6),
                size=Pt(11.5), color=TEXT)
    # right specs
    right_l = Inches(7.4)
    rw = W - right_l - Inches(0.3)
    spec_h = Inches(0.9)
    spec_gap = Inches(0.12)
    for i, (label, val, desc) in enumerate(c['specs']):
        t = Inches(0.9) + i * (spec_h + spec_gap)
        box(sl, right_l, t, rw, spec_h, fill=SURFACE, border_color=BORDER)
        box(sl, right_l, t, Pt(5), spec_h, fill=GREEN)
        tb(sl, label.upper(), right_l + Inches(0.15), t + Inches(0.08),
           rw - Inches(0.2), Inches(0.22), size=Pt(7.5), bold=True, color=MUTED)
        tb(sl, val, right_l + Inches(0.15), t + Inches(0.3),
           rw - Inches(0.2), Inches(0.38), size=Pt(12), bold=True, color=TEXT)
        if desc:
            tb(sl, desc, right_l + Inches(0.15), t + Inches(0.67),
               rw - Inches(0.2), Inches(0.2), size=Pt(8.5), color=MUTED)
    key_bar(sl, c['bar'], accent=BLUE)
    return sl


def slide_architecture(prs, c, n):
    sl = add_slide(prs)
    fade(sl)
    header(sl, c['title'], n=n, total=TOTAL)
    layers = c['layers']
    n_layers = len(layers)
    layer_h = Inches(1.18)
    layer_gap = Inches(0.12)
    arrow_h = Inches(0.18)
    total_h = n_layers * layer_h + (n_layers - 1) * (layer_gap + arrow_h)
    start_t = Inches(0.88)
    arch_w = Inches(12.43)
    arch_l = Inches(0.45)
    cur_t = start_t
    for li, (layer_label, boxes) in enumerate(layers):
        n_boxes = len(boxes)
        box_gap = Inches(0.12)
        box_w = (arch_w - box_gap * (n_boxes - 1)) / n_boxes
        # layer label
        tb(sl, layer_label.upper(), arch_l, cur_t - Inches(0.01), Inches(3), Inches(0.18),
           size=Pt(7.5), bold=True, color=MUTED)
        cur_t += Inches(0.2)
        bh = layer_h - Inches(0.2)
        for bi, (name, sub, col) in enumerate(boxes):
            bl = arch_l + bi * (box_w + box_gap)
            border_c = col if col != MUTED else BORDER
            box(sl, bl, cur_t, box_w, bh, fill=SURFACE, border_color=border_c)
            box(sl, bl, cur_t, box_w, Pt(4), fill=col if col != MUTED else BORDER)
            tb(sl, name, bl + Inches(0.1), cur_t + Inches(0.1),
               box_w - Inches(0.15), Inches(0.3), size=Pt(10), bold=True, color=col if col != MUTED else TEXT)
            tb(sl, sub, bl + Inches(0.1), cur_t + Inches(0.38),
               box_w - Inches(0.15), Inches(0.38), size=Pt(8.5), color=MUTED, wrap=True)
        cur_t += bh
        if li < n_layers - 1:
            tb(sl, '↓', arch_l + arch_w / 2 - Inches(0.2), cur_t + Inches(0.01),
               Inches(0.4), arrow_h, size=Pt(13), color=BORDER, align=PP_ALIGN.CENTER)
            cur_t += layer_gap + arrow_h
    key_bar(sl, c['bar'], accent=BLUE)
    return sl


def slide_ml_model(prs, c, n):
    sl = add_slide(prs)
    fade(sl)
    header(sl, c['title'], n=n, total=TOTAL)
    col_w = Inches(6.0)
    # Left
    tb(sl, c['obj_label'].upper(), Inches(0.45), Inches(0.9), col_w, Inches(0.22),
       size=Pt(8.5), bold=True, color=GREEN)
    callout(sl, c['obj'], Inches(0.45), Inches(1.14), col_w, Inches(0.8), accent=GREEN, size=Pt(11))
    tb(sl, c['algo_label'].upper(), Inches(0.45), Inches(2.05), col_w, Inches(0.22),
       size=Pt(8.5), bold=True, color=BLUE)
    bullet_list(sl, c['algo'], Inches(0.45), Inches(2.3), col_w, Inches(1.6),
                size=Pt(10.5), color=TEXT)
    callout(sl, c['insight'], Inches(0.45), Inches(4.02), col_w, Inches(0.7), accent=BLUE, size=Pt(10))
    # Right — feature table
    right_l = Inches(6.7)
    rw = W - right_l - Inches(0.3)
    tb(sl, c['feat_label'].upper(), right_l, Inches(0.9), rw, Inches(0.22),
       size=Pt(8.5), bold=True, color=GREEN)
    col_widths = [Inches(0.45), Inches(3.7), Inches(1.4)]
    rows = []
    for i, feat_row in enumerate(c['features']):
        if i == 0:
            rows.append(list(feat_row))
        else:
            n_str, feat, pct = feat_row
            color = GREEN if i <= 3 else (AMBER if i <= 5 else MUTED)
            rows.append([(n_str, MUTED, False), (feat, TEXT, False), (pct, color, True)])
    htable(sl, rows, right_l, Inches(1.14), rw, row_h=Inches(0.38),
           col_widths=col_widths)
    key_bar(sl, c['bar'], accent=BLUE)
    return sl


def slide_pipeline(prs, c, n):
    sl = add_slide(prs)
    fade(sl)
    header(sl, c['title'], n=n, total=TOTAL)
    col_w = Inches(6.15)
    col_gap = Inches(0.15)
    left_l  = Inches(0.45)
    right_l = left_l + col_w + col_gap
    step_h  = Inches(0.72)
    step_gap = Inches(0.1)
    accent_colors = [MUTED, MUTED, BLUE, MUTED, MUTED]  # step 3 highlighted
    accent_colors_b = [MUTED, GREEN, GREEN]

    def render_steps(steps, l, start_t, colors):
        t = start_t
        for i, (num, title, desc) in enumerate(steps):
            col = colors[i] if i < len(colors) else MUTED
            bdr = col if col != MUTED else BORDER
            box(sl, l, t, col_w, step_h, fill=SURFACE, border_color=bdr)
            # number bubble
            bub_size = Inches(0.3)
            bub_fill = {BLUE: BLUE_BG, GREEN: GREEN_BG}.get(col, SURFACE)
            box(sl, l + Inches(0.1), t + (step_h - bub_size) / 2, bub_size, bub_size,
                fill=bub_fill, border_color=bdr)
            tb(sl, num, l + Inches(0.1), t + (step_h - bub_size) / 2,
               bub_size, bub_size, size=Pt(9.5), bold=True, color=col,
               align=PP_ALIGN.CENTER)
            tb(sl, title, l + Inches(0.5), t + Inches(0.08),
               col_w - Inches(0.58), Inches(0.28), size=Pt(10.5), bold=True, color=TEXT)
            tb(sl, desc, l + Inches(0.5), t + Inches(0.36),
               col_w - Inches(0.58), Inches(0.3), size=Pt(8.5), color=MUTED, wrap=True)
            t += step_h + step_gap
        return t

    section_label(sl, c['col_a_label'], left_l, Inches(0.88), col_w, accent=BLUE)
    end_a = render_steps(c['steps_a'], left_l, Inches(1.14), accent_colors)

    section_label(sl, c['col_b_label'], right_l, Inches(0.88), col_w, accent=GREEN)
    end_b = render_steps(c['steps_b'], right_l, Inches(1.14), accent_colors_b)
    callout(sl, c['watchdog_note'], right_l, end_b + Inches(0.1), col_w, Inches(0.7), accent=GREEN, size=Pt(9.5))
    key_bar(sl, c['bar'])
    return sl


def slide_risk(prs, c, n):
    sl = add_slide(prs)
    fade(sl)
    header(sl, c['title'], n=n, total=TOTAL)
    cw = Inches(6.15)
    ch = Inches(2.2)
    cg = Inches(0.15)
    l0 = Inches(0.45)
    t0 = Inches(0.9)
    for i, (icon, title, desc) in enumerate(c['items']):
        col = i % 2
        row = i // 2
        l = l0 + col * (cw + cg)
        t = t0 + row * (ch + cg)
        box(sl, l, t, cw, ch, fill=SURFACE, border_color=BORDER)
        tb(sl, icon, l + Inches(0.15), t + Inches(0.12), Inches(0.55), Inches(0.5), size=Pt(22))
        tb(sl, title, l + Inches(0.15), t + Inches(0.62), cw - Inches(0.25), Inches(0.3),
           size=Pt(11.5), bold=True, color=TEXT)
        tb(sl, desc, l + Inches(0.15), t + Inches(0.96), cw - Inches(0.25), Inches(1.1),
           size=Pt(10), color=MUTED, wrap=True)
    key_bar(sl, c['bar'])
    return sl


def slide_metrics(prs, c, n):
    sl = add_slide(prs)
    fade(sl)
    header(sl, c['title'], n=n, total=TOTAL)
    cw = Inches(2.95)
    ch = Inches(1.12)
    cg = Inches(0.15)
    total_w = 4 * cw + 3 * cg
    cl = (W - total_w) / 2
    ct = Inches(0.9)
    for i, (label, val, desc, vc) in enumerate(c['cards']):
        card(sl, cl + i * (cw + cg), ct, cw, ch, label, val, desc, vc=vc, vs=Pt(28))
    # Left column
    lw = Inches(6.15)
    tb(sl, c['interp_label'].upper(), Inches(0.45), Inches(2.22), lw, Inches(0.22),
       size=Pt(8.5), bold=True, color=BLUE)
    bullet_list(sl, c['interp'], Inches(0.45), Inches(2.47), lw, Inches(1.6),
                size=Pt(10.5), color=TEXT)
    # Right column — data table
    right_l = Inches(6.75)
    rw = W - right_l - Inches(0.3)
    tb(sl, c['train_label'].upper(), right_l, Inches(2.22), rw, Inches(0.22),
       size=Pt(8.5), bold=True, color=BLUE)
    rows = [['Key', 'Value']]
    for k, v in c['train_rows']:
        rows.append([(k, TEXT, False), (v, GREEN, True)])
    htable(sl, rows, right_l, Inches(2.47), rw, row_h=Inches(0.38),
           col_widths=[Inches(2.8), Inches(2.8)])
    key_bar(sl, c['bar'], accent=c['bar_accent'])
    return sl


def slide_results(prs, c, n):
    sl = add_slide(prs)
    fade(sl)
    header(sl, c['title'], n=n, total=TOTAL)
    # Hero on the left
    hero_w = Inches(5.5)
    hero_l = Inches(0.45)
    box(sl, hero_l, Inches(0.88), hero_w, Inches(3.0), fill=GREEN_BG, border_color=BORDER)
    box(sl, hero_l, Inches(0.88), hero_w, Pt(4), fill=GREEN)
    tb(sl, c['hero_label'].upper(), hero_l + Inches(0.15), Inches(0.98), hero_w - Inches(0.25), Inches(0.25),
       size=Pt(8.5), bold=True, color=GREEN)
    tb(sl, c['hero_value'], hero_l, Inches(1.22), hero_w, Inches(1.45),
       size=Pt(72), bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    tb(sl, c['hero_sub'], hero_l, Inches(2.72), hero_w, Inches(0.35),
       size=Pt(10.5), color=MUTED, align=PP_ALIGN.CENTER)
    # mini cards
    mini_w = (hero_w - Inches(0.2)) / 2
    mini_h = Inches(0.95)
    mini_t = Inches(3.07)
    for i, (lbl, val, desc, vc) in enumerate(c['cards']):
        card(sl, hero_l + i * (mini_w + Inches(0.1)), mini_t, mini_w, mini_h,
             lbl, val, desc, vc=vc, vs=Pt(26))
    # Right — table + context
    right_l = Inches(6.2)
    rw = W - right_l - Inches(0.3)
    htable(sl, c['table'], right_l, Inches(0.88), rw, row_h=Inches(0.38),
           col_widths=[Inches(2.8), Inches(2.8)])
    callout(sl, c['context'], right_l, Inches(3.0), rw, Inches(1.0), accent=AMBER)
    key_bar(sl, c['bar'])
    return sl


def slide_kelly(prs, c, n):
    sl = add_slide(prs)
    fade(sl)
    header(sl, c['title'], n=n, total=TOTAL)
    lw = Inches(6.8)
    col_widths = [Inches(2.5), Inches(1.4), Inches(1.2), Inches(1.55)]
    table_rows = [c['table'][0]]
    for i, row in enumerate(c['table'][1:]):
        table_rows.append(row)
    # Highlight winner row
    htable(sl, table_rows, Inches(0.45), Inches(0.9), lw,
           row_h=Inches(0.42), col_widths=col_widths)
    # Winner highlight overlay (first data row)
    box(sl, Inches(0.45), Inches(0.9) + Inches(0.42), lw, Pt(3), fill=GREEN)

    right_l = Inches(7.45)
    rw = W - right_l - Inches(0.3)
    tb(sl, c['formula_label'].upper(), right_l, Inches(0.9), rw, Inches(0.22),
       size=Pt(8.5), bold=True, color=AMBER)
    box(sl, right_l, Inches(1.14), rw, Inches(1.3), fill=AMBER_BG, border_color=BORDER)
    box(sl, right_l, Inches(1.14), Pt(5), Inches(1.3), fill=AMBER)
    tb(sl, c['formula'], right_l + Inches(0.15), Inches(1.22), rw - Inches(0.22), Inches(0.32),
       size=Pt(13), bold=True, color=AMBER)
    tb(sl, c['calc'], right_l + Inches(0.15), Inches(1.58), rw - Inches(0.22), Inches(0.32),
       size=Pt(11), bold=True, color=RED)
    tb(sl, c['why'], right_l + Inches(0.15), Inches(2.58), rw - Inches(0.22), Inches(1.0),
       size=Pt(10.5), color=TEXT, wrap=True)
    callout(sl, c['when'], right_l, Inches(3.7), rw, Inches(0.65), accent=GREEN, size=Pt(10))
    key_bar(sl, c['bar'], accent=c['bar_accent'])
    return sl


def slide_dashboard(prs, c, n):
    sl = add_slide(prs)
    fade(sl)
    header(sl, c['title'], n=n, total=TOTAL)
    col_w = Inches(6.7)
    bullet_list(sl, c['features'], Inches(0.45), Inches(0.9), col_w, Inches(3.6),
                size=Pt(11), color=TEXT)
    right_l = Inches(7.35)
    rw = W - right_l - Inches(0.3)
    # Polling intervals
    tb(sl, 'Polling'.upper(), right_l, Inches(0.9), rw, Inches(0.22),
       size=Pt(8.5), bold=True, color=MUTED)
    for i, (interval, desc) in enumerate(c['intervals']):
        t = Inches(1.14) + i * Inches(0.5)
        box(sl, right_l, t, rw, Inches(0.44), fill=SURFACE, border_color=BORDER)
        tb(sl, interval, right_l + Inches(0.1), t + Inches(0.08), Inches(0.55), Inches(0.28),
           size=Pt(12), bold=True, color=BLUE)
        tb(sl, desc, right_l + Inches(0.7), t + Inches(0.1), rw - Inches(0.8), Inches(0.28),
           size=Pt(10), color=TEXT)
    # REST API
    tb(sl, 'REST API', right_l, Inches(2.28), rw, Inches(0.22),
       size=Pt(8.5), bold=True, color=MUTED)
    for i, (method, path, mc) in enumerate(c['api']):
        t = Inches(2.52) + i * Inches(0.42)
        box(sl, right_l, t, rw, Inches(0.36), fill=SURFACE, border_color=BORDER)
        tb(sl, method, right_l + Inches(0.1), t + Inches(0.06), Inches(0.65), Inches(0.26),
           size=Pt(9.5), bold=True, color=mc)
        tb(sl, path, right_l + Inches(0.78), t + Inches(0.07), rw - Inches(0.88), Inches(0.26),
           size=Pt(10), color=TEXT)
    callout(sl, c['note'], right_l, Inches(4.28), rw, Inches(0.65), accent=BLUE, size=Pt(9.5))
    key_bar(sl, c['bar'], accent=BLUE)
    return sl


def slide_team(prs, c, n):
    sl = add_slide(prs)
    fade(sl)
    header(sl, c['title'], n=n, total=TOTAL)
    col_w = Inches(6.0)
    tb(sl, c['process_label'].upper(), Inches(0.45), Inches(0.9), col_w, Inches(0.22),
       size=Pt(8.5), bold=True, color=BLUE)
    bullet_list(sl, c['process'], Inches(0.45), Inches(1.15), col_w, Inches(1.8),
                size=Pt(10.5), color=TEXT)
    callout(sl, c['rule'], Inches(0.45), Inches(3.08), col_w, Inches(0.72), accent=GREEN, size=Pt(10))
    # Team cards
    right_l = Inches(6.65)
    rw = W - right_l - Inches(0.3)
    tb(sl, c['team_label'].upper(), right_l, Inches(0.9), rw, Inches(0.22),
       size=Pt(8.5), bold=True, color=BLUE)
    mh = Inches(0.78)
    mg = Inches(0.1)
    for i, (name, role) in enumerate(c['members']):
        t = Inches(1.15) + i * (mh + mg)
        box(sl, right_l, t, rw, mh, fill=SURFACE, border_color=BORDER)
        box(sl, right_l, t, Pt(5), mh, fill=GREEN)
        tb(sl, name, right_l + Inches(0.15), t + Inches(0.1), rw - Inches(0.2), Inches(0.28),
           size=Pt(11), bold=True, color=TEXT)
        tb(sl, role, right_l + Inches(0.15), t + Inches(0.4), rw - Inches(0.2), Inches(0.32),
           size=Pt(9), color=MUTED, wrap=True)
    key_bar(sl, c['bar'], accent=BLUE)
    return sl


def slide_conclusions(prs, c, n):
    sl = add_slide(prs)
    fade(sl)
    header(sl, c['title'], n=n, total=TOTAL)
    col_w = Inches(4.0)
    col_gap = Inches(0.15)
    cols = [
        (Inches(0.45), c['wins_label'],   GREEN, c['wins']),
        (Inches(0.45) + col_w + col_gap,   c['limits_label'], RED,   c['limits']),
        (Inches(0.45) + 2*(col_w + col_gap), c['next_label'],  BLUE,  c['next']),
    ]
    for l, label, accent, items in cols:
        tb(sl, label.upper(), l, Inches(0.9), col_w, Inches(0.22),
           size=Pt(8.5), bold=True, color=accent)
        box(sl, l, Inches(1.14), col_w, Pt(3), fill=accent)
        bullet_list(sl, items, l, Inches(1.22), col_w, Inches(2.35),
                    size=Pt(10), color=TEXT, spacing=Pt(5))
    callout(sl, c['summary'], Inches(0.45), Inches(3.65), Inches(12.43), Inches(0.8),
            accent=BLUE)
    key_bar(sl, c['bar'])
    return sl


def slide_closing(prs, c):
    sl = add_slide(prs)
    fade(sl)
    box(sl, 0, 0, W, Inches(0.07), fill=GREEN)
    # Left panel (green)
    pw = Inches(4.8)
    box(sl, 0, 0, pw, H, fill=DARK_G)
    box(sl, 0, 0, pw, Inches(0.07), fill=GREEN)
    tb(sl, c['title'], Inches(0.3), Inches(1.5), pw - Inches(0.4), Inches(1.2),
       size=Pt(52), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    box(sl, Inches(0.6), Inches(2.75), pw - Inches(1.1), Pt(2),
        fill=RGBColor(0x7F, 0xD8, 0x9F))
    tb(sl, c['thanks'], Inches(0.3), Inches(2.88), pw - Inches(0.4), Inches(0.75),
       size=Pt(32), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    tb(sl, c['qa'], Inches(0.3), Inches(3.62), pw - Inches(0.4), Inches(0.4),
       size=Pt(13), color=RGBColor(0xAD, 0xE8, 0xC3), align=PP_ALIGN.CENTER)
    # Right — 4 summary cards
    rl = pw + Inches(0.4)
    rw = W - rl - Inches(0.3)
    cw = (rw - Inches(0.15)) / 2
    ch = Inches(1.0)
    items = [
        c['outperform'],
        c['strategy'],
        c['bh'],
        c['winrate'],
    ]
    colors = [GREEN, RED, RED, GREEN]
    for i, ((val, label), col) in enumerate(zip(items, colors)):
        cc = i % 2
        rr = i // 2
        cl = rl + cc * (cw + Inches(0.15))
        ct = Inches(1.3) + rr * (ch + Inches(0.15))
        box(sl, cl, ct, cw, ch, fill=SURFACE, border_color=BORDER)
        box(sl, cl, ct, cw, Pt(4), fill=col)
        tb(sl, val, cl, ct + Inches(0.08), cw, Inches(0.58),
           size=Pt(32), bold=True, color=col, align=PP_ALIGN.CENTER)
        tb(sl, label, cl, ct + Inches(0.66), cw, Inches(0.28),
           size=Pt(8.5), color=MUTED, align=PP_ALIGN.CENTER)
    tb(sl, c['team'], rl, Inches(3.55), rw, Inches(0.32),
       size=Pt(10.5), bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    tb(sl, c['course'], rl, Inches(3.9), rw, Inches(0.28),
       size=Pt(9.5), color=MUTED, align=PP_ALIGN.CENTER)
    tb(sl, c['disclaimer'], rl, Inches(4.25), rw, Inches(0.42),
       size=Pt(8.5), color=MUTED, align=PP_ALIGN.CENTER, wrap=True)
    tb(sl, 'ISLA Bot', Inches(0.45), H - Inches(0.32), Inches(1.5), Inches(0.28),
       size=Pt(8), color=RGBColor(0x7F, 0xD8, 0x9F), bold=True)
    return sl


# ════════════════════════════════════════════════════════════════════════════
# BUILD
# ════════════════════════════════════════════════════════════════════════════

def build(lang):
    c = CONTENT[lang]
    prs = new_prs()
    slide_cover(prs, c['cover'])
    slide_problem(prs, c['problem'], n=2)
    slide_solution(prs, c['solution'], n=3)
    slide_architecture(prs, c['architecture'], n=4)
    slide_ml_model(prs, c['ml_model'], n=5)
    slide_pipeline(prs, c['pipeline'], n=6)
    slide_risk(prs, c['risk'], n=7)
    slide_metrics(prs, c['metrics'], n=8)
    slide_results(prs, c['results'], n=9)
    slide_dashboard(prs, c['dashboard'], n=10)
    slide_team(prs, c['team'], n=11)
    slide_conclusions(prs, c['conclusions'], n=12)
    slide_closing(prs, c['closing'])
    fname = f'isla_bot_{lang}.pptx'
    prs.save(fname)
    print(f'  Saved {fname}  ({len(prs.slides)} slides)')


if __name__ == '__main__':
    print('Building...')
    build('PT')
    build('EN')
    print('Done.')
